"""File converters for Bu D3eij (documents, images, audio/video).

Each converter imports its heavy dependencies lazily for fast startup. Pure
logic — no GUI; ``convert_file`` is the dispatch entry point.
"""
from __future__ import annotations

import sys
from pathlib import Path

from .formats import (
    AV_EXTS,
    ConversionError,
    IMAGE_EXTS,
    PILLOW_FORMAT,
    compatible_targets,
    unique_path,
)


# --------------------------------------------------------------------------- #
# Converters (each imports its heavy deps lazily for fast startup)
# --------------------------------------------------------------------------- #
def convert_image(src: Path, out: Path, target_ext: str) -> None:
    from PIL import Image

    target_ext = target_ext.lower()
    fmt = PILLOW_FORMAT.get(target_ext)
    if not fmt:
        raise ConversionError(f"Unsupported image target: {target_ext}")
    with Image.open(str(src)) as img:
        # Keep animation when both sides support it (GIF/WEBP/TIFF/APNG);
        # otherwise only the first frame survives the conversion.
        animated = getattr(img, "n_frames", 1) > 1
        if animated and target_ext in ("gif", "webp", "tiff", "png"):
            img.save(str(out), fmt, save_all=True)
            return
        # Formats without an alpha channel need a flat RGB image.
        if target_ext in ("jpg", "jpeg", "bmp", "gif") and img.mode in ("RGBA", "LA", "P"):
            img = img.convert("RGB")
        img.save(str(out), fmt)


def pdf_to_txt(src: Path, out: Path) -> None:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(src)) as pdf:
        for page in pdf.pages:
            parts.append(page.extract_text() or "")
    out.write_text("\n\n".join(parts), encoding="utf-8")


def pdf_to_md(src: Path, out: Path) -> None:
    """PDF -> Markdown. Uses pymupdf4llm (structure-aware); falls back to plain text."""
    try:
        import pymupdf4llm

        # Force the classic text-based engine: the optional "layout" path pulls in
        # 20 MB of ONNX models + onnxruntime that we deliberately don't bundle.
        if getattr(pymupdf4llm, "_use_layout", False):
            pymupdf4llm.use_layout(False)
        md = pymupdf4llm.to_markdown(str(src))
    except Exception as exc:  # noqa: BLE001 - fall back to flat text extraction
        print(f"[pdf_to_md] pymupdf4llm unavailable, using text fallback: {exc!r}",
              file=sys.stderr)
        import pdfplumber

        parts: list[str] = []
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages:
                parts.append(page.extract_text() or "")
        md = "\n\n".join(parts)
    out.write_text(md, encoding="utf-8")


def pdf_to_docx(src: Path, out: Path) -> None:
    from pdf2docx import Converter

    cv = Converter(str(src))
    try:
        cv.convert(str(out))  # all pages, layout-aware
    finally:
        cv.close()


def docx_to_txt(src: Path, out: Path) -> None:
    import docx

    document = docx.Document(str(src))
    out.write_text("\n".join(p.text for p in document.paragraphs), encoding="utf-8")


def docx_to_md(src: Path, out: Path) -> None:
    """DOCX -> Markdown via mammoth (DOCX->HTML) + markdownify (HTML->MD)."""
    import mammoth
    from markdownify import markdownify

    with open(src, "rb") as fh:
        html = mammoth.convert_to_html(fh).value
    md = markdownify(html, heading_style="ATX").strip()
    out.write_text(md + "\n", encoding="utf-8")


def docx_to_pdf(src: Path, out: Path) -> None:
    """High-fidelity via MS Word (COM); falls back to a text-only PDF."""
    try:
        _docx_to_pdf_word(src, out)
    except Exception as word_err:  # noqa: BLE001 - want to fall back on anything
        try:
            _docx_to_pdf_reportlab(src, out)
        except Exception as rl_err:  # noqa: BLE001
            raise ConversionError(
                f"DOCX -> PDF failed.\nWord path: {word_err}\nFallback: {rl_err}"
            ) from rl_err


def _docx_to_pdf_word(src: Path, out: Path) -> None:
    # Word is driven through COM directly (not docx2pdf, which only Quit()s Word
    # on success — any failure leaked a hidden WINWORD.EXE per attempt). COM must
    # be initialised on the current (possibly worker) thread. Word.Application is
    # multi-instance, so Dispatch always gives us our own instance to Quit.
    com_ready = False
    try:
        import pythoncom  # provided by pywin32

        pythoncom.CoInitialize()
        com_ready = True
    except Exception:  # noqa: BLE001
        pass
    word = None
    try:
        import win32com.client

        word = win32com.client.Dispatch("Word.Application")
        try:
            word.DisplayAlerts = 0  # never block on dialogs from a hidden instance
        except Exception:  # noqa: BLE001
            pass
        doc = word.Documents.Open(str(src.resolve()), ReadOnly=True)
        try:
            doc.SaveAs(str(out.resolve()), FileFormat=17)  # 17 = wdFormatPDF
        finally:
            doc.Close(0)  # 0 = wdDoNotSaveChanges
    finally:
        if word is not None:
            try:
                word.Quit()
            except Exception:  # noqa: BLE001
                pass
        if com_ready:
            try:
                import pythoncom

                pythoncom.CoUninitialize()
            except Exception:  # noqa: BLE001
                pass
    if not out.exists():
        raise ConversionError("Word produced no output (is Microsoft Word installed?)")


def _docx_to_pdf_reportlab(src: Path, out: Path) -> None:
    from xml.sax.saxutils import escape

    import docx
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    document = docx.Document(str(src))
    styles = getSampleStyleSheet()
    story = []
    for para in document.paragraphs:
        text = para.text.strip()
        if text:
            story.append(Paragraph(escape(text), styles["Normal"]))
            story.append(Spacer(1, 0.08 * inch))
    if not story:
        story.append(Paragraph("(empty document)", styles["Normal"]))
    doc = SimpleDocTemplate(
        str(out), pagesize=LETTER,
        leftMargin=inch, rightMargin=inch, topMargin=inch, bottomMargin=inch,
    )
    doc.build(story)


def _read_text(src: Path) -> str:
    """Read a text file as UTF-8, falling back to cp1252 for legacy files."""
    try:
        return src.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return src.read_text(encoding="cp1252", errors="replace")


def txt_to_md(src: Path, out: Path) -> None:
    """TXT -> Markdown. Plain text is already valid Markdown, so copy it verbatim."""
    out.write_text(_read_text(src), encoding="utf-8")


def txt_to_docx(src: Path, out: Path) -> None:
    """TXT -> DOCX: one paragraph per line (blank lines preserved)."""
    import docx

    document = docx.Document()
    for line in _read_text(src).splitlines():
        document.add_paragraph(line)
    document.save(str(out))


def txt_to_pdf(src: Path, out: Path) -> None:
    """TXT -> PDF: a simple, text-only layout via reportlab (lines wrap)."""
    from xml.sax.saxutils import escape

    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    styles = getSampleStyleSheet()
    story = []
    for line in _read_text(src).splitlines():
        if line.strip():
            story.append(Paragraph(escape(line), styles["Normal"]))
        else:
            story.append(Spacer(1, 0.12 * inch))
    if not story:
        story.append(Paragraph("(empty file)", styles["Normal"]))
    doc = SimpleDocTemplate(
        str(out), pagesize=LETTER,
        leftMargin=inch, rightMargin=inch, topMargin=inch, bottomMargin=inch,
    )
    doc.build(story)


def _pptx_table_to_md(table) -> str:
    """Render a python-pptx table as a GitHub-flavoured Markdown table."""
    def _cell(text: str) -> str:
        # Escape pipes and collapse newlines so cells can't break the table.
        return text.strip().replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ")

    rows = [[_cell(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header, *body = rows
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join("---" for _ in header) + " |"]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)


def pptx_to_md(src: Path, out: Path) -> None:
    """PPTX -> Markdown: one '## Slide' section per slide (text, bullets, tables, notes)."""
    from pptx import Presentation

    prs = Presentation(str(src))
    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title = (title_shape.text.strip() if title_shape and title_shape.text else "") or f"Slide {i}"
        lines = [f"## {title}", ""]
        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            if shape.has_table:
                md_table = _pptx_table_to_md(shape.table)
                if md_table:
                    lines += [md_table, ""]
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append("  " * (para.level or 0) + f"- {text}")
                if shape.text_frame.text.strip():
                    lines.append("")
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines += ["> **Notes:** " + notes.replace("\n", " "), ""]
        blocks.append("\n".join(lines).rstrip())
    out.write_text("\n\n---\n\n".join(blocks) + "\n", encoding="utf-8")


def pptx_to_txt(src: Path, out: Path) -> None:
    """PPTX -> plain text: all shape text per slide, blank line between slides."""
    from pptx import Presentation

    prs = Presentation(str(src))
    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text.strip() for cell in row.cells))
        slides.append("\n".join(parts))
    out.write_text("\n\n".join(slides), encoding="utf-8")


def pptx_to_pdf(src: Path, out: Path) -> None:
    """High-fidelity via MS PowerPoint (COM); falls back to a text-only PDF."""
    try:
        _pptx_to_pdf_powerpoint(src, out)
    except Exception as ppt_err:  # noqa: BLE001 - want to fall back on anything
        try:
            _pptx_to_pdf_reportlab(src, out)
        except Exception as rl_err:  # noqa: BLE001
            raise ConversionError(
                f"PPTX -> PDF failed.\nPowerPoint path: {ppt_err}\nFallback: {rl_err}"
            ) from rl_err


def _pptx_to_pdf_powerpoint(src: Path, out: Path) -> None:
    # PowerPoint is driven through COM, which must be initialised on the current
    # (possibly worker) thread; PowerPoint also requires absolute paths.
    com_ready = False
    try:
        import pythoncom  # provided by pywin32

        pythoncom.CoInitialize()
        com_ready = True
    except Exception:  # noqa: BLE001
        pass
    powerpoint = None
    owned = False  # PowerPoint is single-instance COM: Dispatch attaches to a
    try:           # running PowerPoint, and Quit() would close the user's open
        import win32com.client  # presentations. Only Quit an instance we started.

        try:
            powerpoint = win32com.client.GetActiveObject("PowerPoint.Application")
        except Exception:  # noqa: BLE001 - not running -> start our own
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            owned = True
        presentation = powerpoint.Presentations.Open(
            str(src.resolve()), ReadOnly=True, WithWindow=False
        )
        try:
            presentation.SaveAs(str(out.resolve()), 32)  # 32 = ppSaveAsPDF
        finally:
            presentation.Close()
    finally:
        if powerpoint is not None and owned:
            try:
                powerpoint.Quit()
            except Exception:  # noqa: BLE001
                pass
        if com_ready:
            try:
                import pythoncom

                pythoncom.CoUninitialize()
            except Exception:  # noqa: BLE001
                pass
    if not out.exists():
        raise ConversionError(
            "PowerPoint produced no output (is Microsoft PowerPoint installed?)"
        )


def _pptx_to_pdf_reportlab(src: Path, out: Path) -> None:
    from xml.sax.saxutils import escape

    from pptx import Presentation
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    prs = Presentation(str(src))
    styles = getSampleStyleSheet()
    slides = list(prs.slides)
    story = []
    for i, slide in enumerate(slides, start=1):
        title_shape = slide.shapes.title
        title = (title_shape.text.strip() if title_shape and title_shape.text else "") or f"Slide {i}"
        story.append(Paragraph(escape(title), styles["Heading2"]))
        story.append(Spacer(1, 0.1 * inch))
        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        story.append(Paragraph(escape(text), styles["Normal"]))
                        story.append(Spacer(1, 0.05 * inch))
        if i < len(slides):
            story.append(PageBreak())
    if not story:
        story.append(Paragraph("(empty presentation)", styles["Normal"]))
    doc = SimpleDocTemplate(
        str(out), pagesize=LETTER,
        leftMargin=inch, rightMargin=inch, topMargin=inch, bottomMargin=inch,
    )
    doc.build(story)


def convert_av(src: Path, out: Path, target_ext: str) -> None:
    import ffmpeg

    from .ffmpeg import ensure_ffmpeg

    ensure_ffmpeg()  # system ffmpeg, or download a pinned static build on demand
    target_ext = target_ext.lower()
    src_ext = src.suffix.lower().lstrip(".")

    opts: dict = {}
    if target_ext == "mp3":
        opts = {"acodec": "libmp3lame", "b:a": "192k"}
    elif target_ext == "wav":
        opts = {"acodec": "pcm_s16le"}
    if src_ext == "mp4":  # drop the video stream when extracting audio
        opts["vn"] = None

    stream = ffmpeg.output(ffmpeg.input(str(src)), str(out), **opts)
    try:
        ffmpeg.run(stream, overwrite_output=True, quiet=True)
    except ffmpeg.Error as exc:  # type: ignore[attr-defined]
        detail = exc.stderr.decode("utf-8", "ignore") if getattr(exc, "stderr", None) else str(exc)
        # The full stderr can be many KB and would distort the status label;
        # keep the tail (where ffmpeg states the actual error) for the message.
        print(f"[convert_av] ffmpeg stderr:\n{detail}", file=sys.stderr)
        tail = [ln.strip() for ln in detail.strip().splitlines() if ln.strip()][-3:]
        short = " · ".join(tail)[:400] if tail else str(exc)[:400]
        raise ConversionError(f"ffmpeg failed: {short}") from exc


def convert_file(src, target_ext: str, out_dir=None) -> Path:
    """Convert `src` to `target_ext`. Saves to `out_dir` (or next to the source).

    Returns the output path. Never overwrites an existing file.
    """
    src = Path(src)
    if not src.is_file():
        raise ConversionError(f"File not found: {src}")
    src_ext = src.suffix.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    if target_ext not in compatible_targets(src_ext):
        raise ConversionError(f"Can't convert .{src_ext or '?'} to .{target_ext}")

    out_parent = src.parent
    if out_dir:
        out_parent = Path(out_dir)
        out_parent.mkdir(parents=True, exist_ok=True)
    out = unique_path(out_parent / f"{src.stem}.{target_ext}")
    if src_ext in IMAGE_EXTS:
        convert_image(src, out, target_ext)
    elif src_ext == "pdf" and target_ext == "docx":
        pdf_to_docx(src, out)
    elif src_ext == "pdf" and target_ext == "txt":
        pdf_to_txt(src, out)
    elif src_ext == "pdf" and target_ext == "md":
        pdf_to_md(src, out)
    elif src_ext == "docx" and target_ext == "pdf":
        docx_to_pdf(src, out)
    elif src_ext == "docx" and target_ext == "txt":
        docx_to_txt(src, out)
    elif src_ext == "docx" and target_ext == "md":
        docx_to_md(src, out)
    elif src_ext == "txt" and target_ext == "pdf":
        txt_to_pdf(src, out)
    elif src_ext == "txt" and target_ext == "docx":
        txt_to_docx(src, out)
    elif src_ext == "txt" and target_ext == "md":
        txt_to_md(src, out)
    elif src_ext == "pptx" and target_ext == "md":
        pptx_to_md(src, out)
    elif src_ext == "pptx" and target_ext == "txt":
        pptx_to_txt(src, out)
    elif src_ext == "pptx" and target_ext == "pdf":
        pptx_to_pdf(src, out)
    elif src_ext in AV_EXTS:
        convert_av(src, out, target_ext)
    else:
        raise ConversionError(f"No converter for .{src_ext} to .{target_ext}")
    return out
