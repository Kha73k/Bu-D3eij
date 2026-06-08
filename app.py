"""Bu D3eij — a simple desktop file converter.

A small CustomTkinter app that converts documents, images and audio/video
files. The conversion logic lives in plain module-level functions so it can be
imported and tested without launching the GUI.
"""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import threading
import traceback
from datetime import datetime
from pathlib import Path
from tkinter import filedialog

import customtkinter as ctk
from tkinterdnd2 import DND_FILES, TkinterDnD


def resource_path(rel: str) -> str:
    """Absolute path to a bundled resource (works in dev and in a PyInstaller exe)."""
    base = getattr(sys, "_MEIPASS", os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base, rel)


# --------------------------------------------------------------------------- #
# Format model
# --------------------------------------------------------------------------- #
IMAGE_EXTS = {"jpg", "jpeg", "png", "webp", "bmp", "gif", "tiff"}
DOC_EXTS = {"pdf", "docx", "txt", "md"}
PRESENTATION_EXTS = {"pptx"}
AV_EXTS = {"mp4", "mp3", "wav"}

# input extension -> list of compatible target extensions
CONVERSIONS: dict[str, list[str]] = {
    "pdf": ["docx", "txt", "md"],
    "docx": ["pdf", "txt", "md"],
    "pptx": ["pdf", "txt", "md"],
    "jpg": ["png", "webp", "bmp", "gif", "tiff"],
    "jpeg": ["png", "webp", "bmp", "gif", "tiff"],
    "png": ["jpg", "webp", "bmp", "gif", "tiff"],
    "webp": ["png", "jpg", "bmp", "gif", "tiff"],
    "bmp": ["png", "jpg", "webp", "gif", "tiff"],
    "gif": ["png", "jpg", "webp", "bmp", "tiff"],
    "tiff": ["png", "jpg", "webp", "bmp", "gif"],
    "mp4": ["mp3", "wav"],
    "mp3": ["wav"],
    "wav": ["mp3"],
}

# extension -> Pillow format name
PILLOW_FORMAT = {
    "jpg": "JPEG", "jpeg": "JPEG", "png": "PNG", "webp": "WEBP",
    "bmp": "BMP", "gif": "GIF", "tiff": "TIFF",
}


class ConversionError(Exception):
    """Raised when a conversion cannot be completed."""


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def detect_format(path) -> str:
    """Return the lower-case extension (without dot) of a path."""
    return Path(path).suffix.lower().lstrip(".")


def compatible_targets(src_ext: str) -> list[str]:
    """Return the list of formats `src_ext` can be converted to."""
    return CONVERSIONS.get(src_ext.lower().lstrip("."), [])


def unique_path(path: Path) -> Path:
    """Return a non-existing path by appending ' (n)' before the suffix."""
    if not path.exists():
        return path
    stem, suffix, parent = path.stem, path.suffix, path.parent
    n = 1
    while True:
        candidate = parent / f"{stem} ({n}){suffix}"
        if not candidate.exists():
            return candidate
        n += 1


def human_size(num: float) -> str:
    """Human-readable file size, e.g. '1.4 MB'."""
    for unit in ("B", "KB", "MB", "GB"):
        if num < 1024 or unit == "GB":
            return f"{int(num)} {unit}" if unit == "B" else f"{num:.1f} {unit}"
        num /= 1024
    return f"{num:.1f} GB"


def category_of(ext: str) -> str:
    """Human label for a file extension's category."""
    ext = ext.lower().lstrip(".")
    if ext in IMAGE_EXTS:
        return "Image"
    if ext in PRESENTATION_EXTS:
        return "Presentation"
    if ext in DOC_EXTS:
        return "Document"
    if ext in AV_EXTS:
        return "Audio / Video"
    return "File"


# --------------------------------------------------------------------------- #
# Recent-conversion history (persisted under %LOCALAPPDATA%\Bu D3eij)
# --------------------------------------------------------------------------- #
APP_DATA_DIR = Path(os.environ.get("LOCALAPPDATA") or Path.home()) / "Bu D3eij"
HISTORY_FILE = APP_DATA_DIR / "history.json"
MAX_HISTORY = 100


def load_history() -> list[dict]:
    """Load saved conversion history (newest first). Never raises."""
    try:
        data = json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
        return data if isinstance(data, list) else []
    except Exception:  # noqa: BLE001 - missing/corrupt file is fine
        return []


def save_history(history: list[dict]) -> None:
    """Persist conversion history. Never raises."""
    try:
        APP_DATA_DIR.mkdir(parents=True, exist_ok=True)
        HISTORY_FILE.write_text(json.dumps(history, indent=2), encoding="utf-8")
    except Exception as exc:  # noqa: BLE001
        print("Could not save history:", exc)


# --------------------------------------------------------------------------- #
# Converters (each imports its heavy deps lazily for fast startup)
# --------------------------------------------------------------------------- #
def convert_image(src: Path, out: Path, target_ext: str) -> None:
    from PIL import Image

    target_ext = target_ext.lower()
    fmt = PILLOW_FORMAT.get(target_ext)
    if not fmt:
        raise ConversionError(f"Unsupported image target: {target_ext}")
    with Image.open(str(src)) as img:
        # Formats without an alpha channel need a flat RGB image.
        if target_ext in ("jpg", "jpeg", "bmp", "gif") and img.mode in ("RGBA", "LA", "P"):
            img = img.convert("RGB")
        img.save(str(out), fmt)


def pdf_to_txt(src: Path, out: Path) -> None:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(src)) as pdf:
        for page in pdf.pages:
            parts.append(page.extract_text() or "")
    out.write_text("\n\n".join(parts), encoding="utf-8")


def pdf_to_md(src: Path, out: Path) -> None:
    """PDF -> Markdown. Uses pymupdf4llm (structure-aware); falls back to plain text."""
    try:
        import pymupdf4llm

        # Force the classic text-based engine: the optional "layout" path pulls in
        # 20 MB of ONNX models + onnxruntime that we deliberately don't bundle.
        if getattr(pymupdf4llm, "_use_layout", False):
            pymupdf4llm.use_layout(False)
        md = pymupdf4llm.to_markdown(str(src))
    except Exception as exc:  # noqa: BLE001 - fall back to flat text extraction
        print(f"[pdf_to_md] pymupdf4llm unavailable, using text fallback: {exc!r}",
              file=sys.stderr)
        import pdfplumber

        parts: list[str] = []
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages:
                parts.append(page.extract_text() or "")
        md = "\n\n".join(parts)
    out.write_text(md, encoding="utf-8")


def pdf_to_docx(src: Path, out: Path) -> None:
    from pdf2docx import Converter

    cv = Converter(str(src))
    try:
        cv.convert(str(out))  # all pages, layout-aware
    finally:
        cv.close()


def docx_to_txt(src: Path, out: Path) -> None:
    import docx

    document = docx.Document(str(src))
    out.write_text("\n".join(p.text for p in document.paragraphs), encoding="utf-8")


def docx_to_md(src: Path, out: Path) -> None:
    """DOCX -> Markdown via mammoth (DOCX->HTML) + markdownify (HTML->MD)."""
    import mammoth
    from markdownify import markdownify

    with open(src, "rb") as fh:
        html = mammoth.convert_to_html(fh).value
    md = markdownify(html, heading_style="ATX").strip()
    out.write_text(md + "\n", encoding="utf-8")


def docx_to_pdf(src: Path, out: Path) -> None:
    """High-fidelity via MS Word (docx2pdf); falls back to a text-only PDF."""
    try:
        _docx_to_pdf_word(src, out)
    except Exception as word_err:  # noqa: BLE001 - want to fall back on anything
        try:
            _docx_to_pdf_reportlab(src, out)
        except Exception as rl_err:  # noqa: BLE001
            raise ConversionError(
                f"DOCX -> PDF failed.\nWord path: {word_err}\nFallback: {rl_err}"
            ) from rl_err


def _docx_to_pdf_word(src: Path, out: Path) -> None:
    # docx2pdf drives MS Word through COM, which must be initialised on the
    # current (possibly worker) thread.
    com_ready = False
    try:
        import pythoncom  # provided by pywin32

        pythoncom.CoInitialize()
        com_ready = True
    except Exception:  # noqa: BLE001
        pass
    try:
        from docx2pdf import convert as docx2pdf_convert

        docx2pdf_convert(str(src), str(out))
    finally:
        if com_ready:
            try:
                import pythoncom

                pythoncom.CoUninitialize()
            except Exception:  # noqa: BLE001
                pass
    if not out.exists():
        raise ConversionError("docx2pdf produced no output (is Microsoft Word installed?)")


def _docx_to_pdf_reportlab(src: Path, out: Path) -> None:
    from xml.sax.saxutils import escape

    import docx
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    document = docx.Document(str(src))
    styles = getSampleStyleSheet()
    story = []
    for para in document.paragraphs:
        text = para.text.strip()
        if text:
            story.append(Paragraph(escape(text), styles["Normal"]))
            story.append(Spacer(1, 0.08 * inch))
    if not story:
        story.append(Paragraph("(empty document)", styles["Normal"]))
    doc = SimpleDocTemplate(
        str(out), pagesize=LETTER,
        leftMargin=inch, rightMargin=inch, topMargin=inch, bottomMargin=inch,
    )
    doc.build(story)


def _pptx_table_to_md(table) -> str:
    """Render a python-pptx table as a GitHub-flavoured Markdown table."""
    def _cell(text: str) -> str:
        # Escape pipes and collapse newlines so cells can't break the table.
        return text.strip().replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ")

    rows = [[_cell(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header, *body = rows
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join("---" for _ in header) + " |"]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)


def pptx_to_md(src: Path, out: Path) -> None:
    """PPTX -> Markdown: one '## Slide' section per slide (text, bullets, tables, notes)."""
    from pptx import Presentation

    prs = Presentation(str(src))
    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title = (title_shape.text.strip() if title_shape and title_shape.text else "") or f"Slide {i}"
        lines = [f"## {title}", ""]
        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            if shape.has_table:
                md_table = _pptx_table_to_md(shape.table)
                if md_table:
                    lines += [md_table, ""]
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append("  " * (para.level or 0) + f"- {text}")
                if shape.text_frame.text.strip():
                    lines.append("")
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines += ["> **Notes:** " + notes.replace("\n", " "), ""]
        blocks.append("\n".join(lines).rstrip())
    out.write_text("\n\n---\n\n".join(blocks) + "\n", encoding="utf-8")


def pptx_to_txt(src: Path, out: Path) -> None:
    """PPTX -> plain text: all shape text per slide, blank line between slides."""
    from pptx import Presentation

    prs = Presentation(str(src))
    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text.strip() for cell in row.cells))
        slides.append("\n".join(parts))
    out.write_text("\n\n".join(slides), encoding="utf-8")


def pptx_to_pdf(src: Path, out: Path) -> None:
    """High-fidelity via MS PowerPoint (COM); falls back to a text-only PDF."""
    try:
        _pptx_to_pdf_powerpoint(src, out)
    except Exception as ppt_err:  # noqa: BLE001 - want to fall back on anything
        try:
            _pptx_to_pdf_reportlab(src, out)
        except Exception as rl_err:  # noqa: BLE001
            raise ConversionError(
                f"PPTX -> PDF failed.\nPowerPoint path: {ppt_err}\nFallback: {rl_err}"
            ) from rl_err


def _pptx_to_pdf_powerpoint(src: Path, out: Path) -> None:
    # PowerPoint is driven through COM, which must be initialised on the current
    # (possibly worker) thread; PowerPoint also requires absolute paths.
    com_ready = False
    try:
        import pythoncom  # provided by pywin32

        pythoncom.CoInitialize()
        com_ready = True
    except Exception:  # noqa: BLE001
        pass
    powerpoint = None
    try:
        import win32com.client

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(src.resolve()), ReadOnly=True, WithWindow=False
        )
        try:
            presentation.SaveAs(str(out.resolve()), 32)  # 32 = ppSaveAsPDF
        finally:
            presentation.Close()
    finally:
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:  # noqa: BLE001
                pass
        if com_ready:
            try:
                import pythoncom

                pythoncom.CoUninitialize()
            except Exception:  # noqa: BLE001
                pass
    if not out.exists():
        raise ConversionError(
            "PowerPoint produced no output (is Microsoft PowerPoint installed?)"
        )


def _pptx_to_pdf_reportlab(src: Path, out: Path) -> None:
    from xml.sax.saxutils import escape

    from pptx import Presentation
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    prs = Presentation(str(src))
    styles = getSampleStyleSheet()
    slides = list(prs.slides)
    story = []
    for i, slide in enumerate(slides, start=1):
        title_shape = slide.shapes.title
        title = (title_shape.text.strip() if title_shape and title_shape.text else "") or f"Slide {i}"
        story.append(Paragraph(escape(title), styles["Heading2"]))
        story.append(Spacer(1, 0.1 * inch))
        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        story.append(Paragraph(escape(text), styles["Normal"]))
                        story.append(Spacer(1, 0.05 * inch))
        if i < len(slides):
            story.append(PageBreak())
    if not story:
        story.append(Paragraph("(empty presentation)", styles["Normal"]))
    doc = SimpleDocTemplate(
        str(out), pagesize=LETTER,
        leftMargin=inch, rightMargin=inch, topMargin=inch, bottomMargin=inch,
    )
    doc.build(story)


def convert_av(src: Path, out: Path, target_ext: str) -> None:
    import ffmpeg

    if shutil.which("ffmpeg") is None:
        raise ConversionError(
            "ffmpeg is not installed or not on PATH. Install it "
            "(e.g. 'winget install Gyan.FFmpeg') and restart the app."
        )
    target_ext = target_ext.lower()
    src_ext = src.suffix.lower().lstrip(".")

    opts: dict = {}
    if target_ext == "mp3":
        opts = {"acodec": "libmp3lame", "b:a": "192k"}
    elif target_ext == "wav":
        opts = {"acodec": "pcm_s16le"}
    if src_ext == "mp4":  # drop the video stream when extracting audio
        opts["vn"] = None

    stream = ffmpeg.output(ffmpeg.input(str(src)), str(out), **opts)
    try:
        ffmpeg.run(stream, overwrite_output=True, quiet=True)
    except ffmpeg.Error as exc:  # type: ignore[attr-defined]
        detail = exc.stderr.decode("utf-8", "ignore") if getattr(exc, "stderr", None) else str(exc)
        raise ConversionError(f"ffmpeg failed: {detail}") from exc


def convert_file(src, target_ext: str, out_dir=None) -> Path:
    """Convert `src` to `target_ext`. Saves to `out_dir` (or next to the source).

    Returns the output path. Never overwrites an existing file.
    """
    src = Path(src)
    if not src.is_file():
        raise ConversionError(f"File not found: {src}")
    src_ext = src.suffix.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    if target_ext not in compatible_targets(src_ext):
        raise ConversionError(f"Can't convert .{src_ext or '?'} to .{target_ext}")

    out_parent = src.parent
    if out_dir:
        out_parent = Path(out_dir)
        out_parent.mkdir(parents=True, exist_ok=True)
    out = unique_path(out_parent / f"{src.stem}.{target_ext}")
    if src_ext in IMAGE_EXTS:
        convert_image(src, out, target_ext)
    elif src_ext == "pdf" and target_ext == "docx":
        pdf_to_docx(src, out)
    elif src_ext == "pdf" and target_ext == "txt":
        pdf_to_txt(src, out)
    elif src_ext == "pdf" and target_ext == "md":
        pdf_to_md(src, out)
    elif src_ext == "docx" and target_ext == "pdf":
        docx_to_pdf(src, out)
    elif src_ext == "docx" and target_ext == "txt":
        docx_to_txt(src, out)
    elif src_ext == "docx" and target_ext == "md":
        docx_to_md(src, out)
    elif src_ext == "pptx" and target_ext == "md":
        pptx_to_md(src, out)
    elif src_ext == "pptx" and target_ext == "txt":
        pptx_to_txt(src, out)
    elif src_ext == "pptx" and target_ext == "pdf":
        pptx_to_pdf(src, out)
    elif src_ext in AV_EXTS:
        convert_av(src, out, target_ext)
    else:
        raise ConversionError(f"No converter for .{src_ext} to .{target_ext}")
    return out


# --------------------------------------------------------------------------- #
# YouTube download (yt-dlp + ffmpeg)
# --------------------------------------------------------------------------- #
YT_FORMATS = ("mp3", "mp4")


def download_youtube(url: str, fmt: str, out_dir, progress_hook=None) -> Path:
    """Download a YouTube (or other yt-dlp-supported) URL as mp3 or mp4.

    `fmt` is "mp3" (192 kbps audio) or "mp4" (best video+audio, merged).
    Saves into `out_dir` and returns the output path. Needs ffmpeg.
    """
    import yt_dlp

    fmt = fmt.lower().lstrip(".")
    if fmt not in YT_FORMATS:
        raise ConversionError(f"Unsupported download format: {fmt}")
    url = (url or "").strip()
    if not url.lower().startswith(("http://", "https://")):
        raise ConversionError("Enter a valid video URL (must start with http:// or https://).")
    if shutil.which("ffmpeg") is None:
        raise ConversionError(
            "ffmpeg is not installed or not on PATH. Install it "
            "(e.g. 'winget install Gyan.FFmpeg') and restart the app."
        )

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    outtmpl = str(out_dir / "%(title)s.%(ext)s")
    opts: dict = {
        "outtmpl": outtmpl,
        "ffmpeg_location": shutil.which("ffmpeg"),
        "quiet": True,
        "no_warnings": True,
        "noplaylist": True,  # a single video, even if the URL has a list= param
        "restrictfilenames": True,  # filesystem-safe names
    }
    if progress_hook:
        opts["progress_hooks"] = [progress_hook]
    if fmt == "mp3":
        opts["format"] = "bestaudio/best"
        opts["postprocessors"] = [{
            "key": "FFmpegExtractAudio",
            "preferredcodec": "mp3",
            "preferredquality": "192",
        }]
    else:  # mp4
        opts["format"] = "bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]/best"
        opts["merge_output_format"] = "mp4"

    try:
        with yt_dlp.YoutubeDL(opts) as ydl:
            info = ydl.extract_info(url, download=True)
            produced = Path(ydl.prepare_filename(info))
    except yt_dlp.utils.DownloadError as exc:
        raise ConversionError(f"Download failed: {exc}") from exc
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"Download failed: {exc}") from exc

    # The post-processor changes the extension; resolve the final file.
    out = produced.with_suffix(f".{fmt}")
    if not out.exists():
        # Fall back to whatever landed if the suffix guess is off.
        if produced.exists():
            return produced
        raise ConversionError("Download finished but the output file was not found.")
    return out


# Lazily-created rembg sessions, cached per model name (model loaded once each).
_REMBG_SESSIONS: dict = {}


def remove_background(src, out_path=None, model: str = "isnet-general-use") -> Path:
    """Remove an image's background and save a transparent PNG; return its path.

    `model` is a rembg model name (see `BG_MODELS` for the Marquee tiers:
    `u2netp` / `isnet-general-use` / `birefnet-general`). `src` must be an image;
    the output is always saved as PNG (the only supported image format that
    preserves transparency). Like `download_youtube`, this lives outside
    `convert_file`/`CONVERSIONS` — it has no target format to choose.
    """
    from PIL import Image
    from rembg import new_session, remove  # lazy: heavy (onnxruntime + model)

    src = Path(src)
    ext = src.suffix.lower().lstrip(".")
    if ext not in IMAGE_EXTS:
        raise ConversionError(f"Background removal needs an image; got .{ext or '?'}")

    try:
        session = _REMBG_SESSIONS.get(model)
        if session is None:
            # Downloads this model's .onnx on first use, then caches in ~/.u2net.
            session = new_session(model)
            _REMBG_SESSIONS[model] = session
        with Image.open(src) as im:
            result = remove(im.convert("RGBA"), session=session)
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"Background removal failed: {exc}") from exc

    out = Path(out_path) if out_path else src.with_name(f"{src.stem}_no-bg.png")
    if out.suffix.lower() != ".png":
        out = out.with_suffix(".png")  # transparency requires PNG
    out = unique_path(out)             # never clobber an existing file
    result.save(out, "PNG")
    return out


# --------------------------------------------------------------------------- #
# GUI
# --------------------------------------------------------------------------- #
APP_NAME = "Bu D3eij"
NAV_ITEMS = ["Home", "Converter", "Recent", "Batch Convert", "YouTube", "Marquee", "Tools"]

# Logo-derived palette (extracted from AppLogo.png).
RED = "#E11414"
RED_HOVER = "#B4000C"
RED_DEEP = "#8C0008"
RED_BRIGHT = "#F01818"
SIDEBAR_FG = ("#F3EEEF", "#141011")
NAV_TEXT = ("#2A2426", "#F2E9EA")
DROP_BORDER = ("#E11414", "#B4000C")
SUCCESS = "#1FA85B"
ERROR = "#F0282D"
MUTED = ("#6B6164", "#9C9194")
# Surfaces for the 1.4 redesign (light, dark).
CARD = ("#FFFFFF", "#252022")          # elevated card / panel
CARD_BORDER = ("#E8E1E2", "#332D30")   # hairline card border
SURFACE_SOFT = ("#FBECEC", "#221A1B")  # subtle red-tinted hero / accent surface
TEXT = ("#1A1416", "#F2E9EA")          # primary text
SUN_GLYPH = "☀"   # ☀ shown while in Dark mode (click → Light)
MOON_GLYPH = "☾"  # ☾ shown while in Light mode (click → Dark)
APP_VERSION = "2.1"

# Extension -> file-type icon (assets/filetypes/<key>.png). Falls back to "default".
EXT_ICON = {
    "pdf": "pdf", "docx": "word", "doc": "word", "txt": "text", "md": "markdown",
    "pptx": "powerpoint", "ppt": "powerpoint",
    "jpg": "image", "jpeg": "image", "png": "image", "webp": "image",
    "bmp": "image", "gif": "image", "tiff": "image",
    "mp4": "video", "mp3": "audio", "wav": "audio",
}


def icon_key_for_ext(ext: str) -> str:
    return EXT_ICON.get(ext.lower().lstrip("."), "default")


# Nav label -> UI icon (assets/ui/<name>.png).
NAV_ICONS = {
    "Home": "house", "Converter": "repeat", "Recent": "clock",
    "Batch Convert": "layers", "YouTube": "youtube", "Marquee": "sparkles",
    "Tools": "wrench",
}

# Marquee background-removal tiers -> (rembg model name, one-line blurb).
# All three ride on the bundled rembg; each downloads its own .onnx on first use.
BG_MODELS = {
    "Flash": ("u2netp",            "Fastest · lightweight model for quick cut-outs"),
    "Mid":   ("isnet-general-use", "Balanced · cleaner edges at near-Flash speed"),
    "Omega": ("birefnet-general",  "Max precision · best hair/edges, slower + larger download"),
}
DEFAULT_BG_TIER = "Mid"

ctk.set_appearance_mode("Dark")
try:
    ctk.set_default_color_theme(resource_path("bud3eij_theme.json"))
except Exception:  # noqa: BLE001 - fall back to a built-in theme
    ctk.set_default_color_theme("blue")


class GradientButton(ctk.CTkFrame):
    """A flashy, fully animated 'Convert' button drawn entirely with Pillow.

    CustomTkinter (Tkinter) has no CSS engine — no gradients, sweeping shines,
    glows or transitions — so the button's whole visual is composed as a PIL
    image and swapped onto an inner label every animation tick. Effects:
      * red left→right gradient with a top gloss sheen,
      * a light band that sweeps across (the "shine"),
      * a breathing glow when idle, a brighter glow + faster shine on hover,
      * a quick darken on press,
      * a continuous double-shine "flow" with animated "Converting…" dots while
        a conversion runs (``start_busy`` / ``stop_busy``),
      * a flat, greyed, motionless look while disabled.

    Drop-in for the old ``CTkButton``: supports ``grid``, a ``command`` callback
    and ``configure(state="normal"|"disabled")``. Animation only runs while the
    widget is mapped (it pauses when you switch pages) and is fully cancelled on
    destroy, so it never burns CPU in the background.
    """

    GLOW_PAD = 9      # logical px of glow margin reserved around the body
    RADIUS = 0.34     # corner radius as a fraction of the body height

    def __init__(self, master, text: str = "Convert", height: int = 46,
                 command=None, state: str = "normal", **_ignored):
        super().__init__(master, fg_color="transparent",
                         height=height + 2 * self.GLOW_PAD)
        self._text = text
        self._command = command
        self._btn_h = height
        self._enabled = state != "disabled"
        self._busy = False
        self._hover = False
        self._press = False
        # animation state
        self._phase = 0.0     # shine sweep position 0..1
        self._breath = 0.0    # breathing oscillator (radians)
        self._glow = 0.0      # eased glow amount 0..1
        self._dots = 0        # animated "Converting…" dot count
        self._dot_acc = 0
        self._anim_id = None
        self._resize_id = None
        # cached static layers (rebuilt on size/state change)
        self._stat = None
        self._stat_key = None
        self._cur_img = None  # keep a ref so the PhotoImage isn't GC'd

        self.grid_propagate(False)
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(0, weight=1)
        self._label = ctk.CTkLabel(self, text="", fg_color="transparent")
        self._label.grid(row=0, column=0, sticky="nsew")

        for w in (self, self._label):
            w.bind("<Enter>", self._on_enter)
            w.bind("<Leave>", self._on_leave)
            w.bind("<Button-1>", self._on_press)
            w.bind("<ButtonRelease-1>", self._on_release)
        self.bind("<Configure>", self._on_configure)
        self.bind("<Map>", lambda _e: self._start())
        self.bind("<Unmap>", lambda _e: self._stop())
        self.bind("<Destroy>", self._on_destroy)

    # ---- public drop-in API --------------------------------------------- #
    def configure(self, **kwargs):
        if "state" in kwargs:
            self.set_enabled(kwargs.pop("state") != "disabled")
        if "text" in kwargs:
            self._text = kwargs.pop("text")
            self._render()
        if kwargs:
            super().configure(**kwargs)

    def set_enabled(self, on: bool):
        on = bool(on)
        if on == self._enabled:
            return
        self._enabled = on
        if not on:
            self._hover = self._press = False
        if on:
            self._start()
        elif not self._busy:
            self._stop()
        self._render()

    def start_busy(self):
        """Switch to the animated 'Converting…' look (clicks ignored)."""
        self._busy = True
        self._phase = 0.0
        self._dots = 0
        self._dot_acc = 0
        self._start()

    def stop_busy(self):
        self._busy = False
        if self._enabled:
            self._start()
        else:
            self._stop()
        self._render()

    # ---- event handlers -------------------------------------------------- #
    def _on_enter(self, _e):
        if self._enabled and not self._busy:
            self._hover = True

    def _on_leave(self, _e):
        self._hover = False
        self._press = False

    def _on_press(self, _e):
        if self._enabled and not self._busy:
            self._press = True

    def _on_release(self, _e):
        fire = self._press and self._hover and self._enabled and not self._busy
        self._press = False
        if fire and self._command is not None:
            self._command()

    def _on_configure(self, _e):
        if self._resize_id is not None:
            try:
                self.after_cancel(self._resize_id)
            except Exception:  # noqa: BLE001
                pass
        self._resize_id = self.after(60, self._apply_resize)

    def _apply_resize(self):
        self._resize_id = None
        self._stat = None  # force static layers to rebuild at the new width
        self._render()

    def _on_destroy(self, event):
        if event.widget is self:
            self._stop()

    # ---- animation loop -------------------------------------------------- #
    def _start(self):
        if self._anim_id is not None:
            return
        try:
            if not self.winfo_ismapped():
                return
        except Exception:  # noqa: BLE001
            return
        if not (self._busy or self._enabled):
            self._render()
            return
        self._tick()

    def _stop(self):
        if self._anim_id is not None:
            try:
                self.after_cancel(self._anim_id)
            except Exception:  # noqa: BLE001
                pass
            self._anim_id = None

    def _tick(self):
        import math
        self._anim_id = None
        if not self.winfo_exists() or not self.winfo_ismapped():
            return
        if self._busy:
            self._phase = (self._phase + 0.035) % 1.0
            target_glow = 0.32 + 0.12 * (0.5 + 0.5 * math.sin(self._breath))
            interval = 33
        elif self._hover:
            self._phase = (self._phase + 0.022) % 1.0
            target_glow = 0.60
            interval = 33
        else:
            self._phase = (self._phase + 0.012) % 1.0
            target_glow = 0.10 + 0.08 * (0.5 + 0.5 * math.sin(self._breath))
            interval = 50
        if self._press:
            target_glow = 0.25
        self._breath += 0.10
        self._glow += (target_glow - self._glow) * 0.25
        if self._busy:
            self._dot_acc += interval
            if self._dot_acc >= 400:
                self._dot_acc = 0
                self._dots = (self._dots + 1) % 4
        self._render()
        if self._busy or self._enabled:
            self._anim_id = self.after(interval, self._tick)

    # ---- rendering ------------------------------------------------------- #
    @staticmethod
    def _load_font(px: int):
        from PIL import ImageFont
        px = max(8, px)
        fonts_dir = os.path.join(os.environ.get("WINDIR", "C:\\Windows"), "Fonts")
        for name in ("segoeuib.ttf", "seguisb.ttf", "arialbd.ttf"):
            try:
                return ImageFont.truetype(os.path.join(fonts_dir, name), px)
            except Exception:  # noqa: BLE001
                continue
        try:
            return ImageFont.truetype("arialbd.ttf", px)
        except Exception:  # noqa: BLE001
            return ImageFont.load_default()

    def _load_spark(self, px: int):
        try:
            from PIL import Image

            s = Image.open(resource_path("assets/ui/sparkles.png")).convert("RGBA")
            white = Image.new("RGBA", s.size, (255, 255, 255, 0))
            white.putalpha(s.split()[3])
            scale = px / max(white.size)
            return white.resize((max(1, int(white.width * scale)),
                                 max(1, int(white.height * scale))))
        except Exception:  # noqa: BLE001
            return None

    def _ensure_static(self, w: int, h: int, scaling: float):
        key = (w, h, self._enabled)
        if self._stat_key == key and self._stat is not None:
            return self._stat
        import math

        from PIL import Image, ImageDraw, ImageFilter

        pad = max(1, int(round(self.GLOW_PAD * scaling)))
        rw = max(1, w - 2 * pad)
        rh = max(1, h - 2 * pad)
        radius = max(1, int(rh * self.RADIUS))

        mask = Image.new("L", (rw, rh), 0)
        ImageDraw.Draw(mask).rounded_rectangle([0, 0, rw - 1, rh - 1], radius, fill=255)

        if self._enabled:
            c_l, c_r = (0xF0, 0x2A, 0x2A), (0x8C, 0x00, 0x08)
            gloss_peak = 60
        else:
            c_l, c_r = (0x6E, 0x66, 0x68), (0x4A, 0x44, 0x46)
            gloss_peak = 28
        row = Image.new("RGB", (rw, 1))
        rpx = row.load()
        for x in range(rw):
            t = x / max(1, rw - 1)
            rpx[x, 0] = (int(c_l[0] + (c_r[0] - c_l[0]) * t),
                         int(c_l[1] + (c_r[1] - c_l[1]) * t),
                         int(c_l[2] + (c_r[2] - c_l[2]) * t))
        base = row.resize((rw, rh))

        gcol = Image.new("L", (1, rh))
        gpx = gcol.load()
        for y in range(rh):
            ty = y / max(1, rh - 1)
            gpx[0, y] = int(max(0.0, (0.5 - ty) / 0.5) * gloss_peak)
        base = Image.composite(Image.new("RGB", (rw, rh), (255, 255, 255)),
                               base, gcol.resize((rw, rh)))

        band_w = max(int(rh * 1.4), int(rw * 0.15), 2)
        sb = Image.new("L", (band_w, 1))
        spx = sb.load()
        sigma = max(1.0, band_w * 0.22)
        for x in range(band_w):
            d = (x - band_w / 2) / sigma
            spx[x, 0] = int(150 * math.exp(-0.5 * d * d))
        shine_a = sb.resize((band_w, rh))
        shine_w = Image.new("RGB", (band_w, rh), (255, 255, 255))

        glow = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        if self._enabled:
            ImageDraw.Draw(glow).rounded_rectangle(
                [pad - 2, pad - 2, w - pad + 1, h - pad + 1], radius + 2,
                fill=(0xFF, 0x22, 0x22, 255))
            glow = glow.filter(ImageFilter.GaussianBlur(max(2, int(pad * 1.1))))

        self._stat = dict(
            pad=pad, rw=rw, rh=rh, mask=mask, base=base, band_w=band_w,
            shine_a=shine_a, shine_w=shine_w, glow=glow,
            font=self._load_font(int(round(self._btn_h * 0.36 * scaling))),
            spark=self._load_spark(int(round(self._btn_h * 0.42 * scaling))),
        )
        self._stat_key = key
        return self._stat

    def _compose(self):
        import math

        from PIL import Image, ImageDraw, ImageEnhance

        w = self.winfo_width()
        h = self.winfo_height()
        if w <= 1 or h <= 1:
            return None
        scaling = self._get_widget_scaling()
        st = self._ensure_static(w, h, scaling)
        pad, rw, rh = st["pad"], st["rw"], st["rh"]

        if not self._enabled and not self._busy:
            factor = 1.0
        elif self._press:
            factor = 0.90
        elif self._hover:
            factor = 1.14
        elif self._busy:
            factor = 1.06
        else:
            factor = 1.0 + 0.05 * (0.5 + 0.5 * math.sin(self._breath))

        body = st["base"].copy()
        if abs(factor - 1.0) > 1e-3:
            body = ImageEnhance.Brightness(body).enhance(factor)

        if self._enabled or self._busy:
            bw = st["band_w"]
            spans = [self._phase]
            if self._busy:
                spans.append((self._phase + 0.5) % 1.0)
            for ph in spans:
                sx = int(ph * (rw + bw)) - bw
                body.paste(st["shine_w"], (sx, 0), st["shine_a"])

        draw = ImageDraw.Draw(body)
        label = ("Converting" + "." * self._dots) if self._busy else self._text
        txt_color = (255, 255, 255) if (self._enabled or self._busy) else (190, 182, 184)
        font, spark = st["font"], st["spark"]
        bbox = draw.textbbox((0, 0), label, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        gap = max(2, int(rh * 0.12))
        show_spark = spark is not None and (self._enabled or self._busy)
        sw = (spark.width + gap) if show_spark else 0
        x0 = (rw - (tw + sw)) // 2
        y0 = (rh - th) // 2 - bbox[1]
        if show_spark:
            body.paste(spark, (x0, (rh - spark.height) // 2), spark)
            x0 += spark.width + gap
        draw.text((x0, y0), label, font=font, fill=txt_color)

        body = body.convert("RGBA")
        body.putalpha(st["mask"])

        full = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        if self._glow > 0.01 and self._enabled and not self._press:
            g = st["glow"].copy()
            amt = min(1.0, self._glow)
            g.putalpha(g.split()[3].point(lambda v: int(v * amt)))
            full.alpha_composite(g)
        full.alpha_composite(body, (pad, pad))
        return full

    def _render(self):
        if not self.winfo_exists():
            return
        try:
            img = self._compose()
        except Exception as exc:  # noqa: BLE001
            print("convert button render failed:", exc)
            return
        if img is None:
            return
        scaling = self._get_widget_scaling()
        ci = ctk.CTkImage(light_image=img, dark_image=img,
                          size=(img.width / scaling, img.height / scaling))
        self._cur_img = ci
        self._label.configure(image=ci)


class App(ctk.CTk, TkinterDnD.DnDWrapper):
    def __init__(self):
        super().__init__()
        # Enable native drag & drop on a CustomTkinter root.
        self.TkdndVersion = TkinterDnD._require(self)

        self.title(APP_NAME)
        self.geometry("1000x680")
        self.minsize(880, 600)
        self._set_window_icon()

        self.selected_file: Path | None = None
        self.export_dir: Path | None = None
        self.batch_files: list[Path] = []
        self.marquee_file: Path | None = None
        self.history: list[dict] = load_history()
        self.appearance_mode = "Dark"  # toggled by the sun/moon button
        self._icon_cache: dict = {}  # (kind, name, size, colors) -> CTkImage

        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        self._build_sidebar()
        self._build_frames()
        self.show_frame("Converter")

    def _set_window_icon(self):
        self._icon_path = resource_path("AppLogo.ico")
        if not os.path.exists(self._icon_path):
            return
        self._apply_icon()
        # Re-apply shortly after startup (CustomTkinter sets its own icon late).
        self.after(300, self._apply_icon)

    def _apply_icon(self):
        try:
            if self.winfo_exists():
                self.iconbitmap(self._icon_path)
        except Exception as exc:  # noqa: BLE001
            print("Could not set window icon:", exc)

    def _load_logo(self, name: str, width: int):
        """Load a logo as a width-scaled CTkImage, or None on failure."""
        try:
            from PIL import Image

            img = Image.open(resource_path(name))
            w, h = img.size
            return ctk.CTkImage(
                light_image=img, dark_image=img, size=(width, int(width * h / w))
            )
        except Exception as exc:  # noqa: BLE001
            print("Could not load logo", name, ":", exc)
            return None

    # ---- icons ----------------------------------------------------------- #
    @staticmethod
    def _hex_rgb(color: str) -> tuple[int, int, int]:
        c = color.lstrip("#")
        return int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)

    def _filetype_icon(self, ext: str, size: int = 36):
        """Colored file-type badge (CTkImage) for a file extension. Cached."""
        key = ("ft", icon_key_for_ext(ext), size)
        if key in self._icon_cache:
            return self._icon_cache[key]
        try:
            from PIL import Image

            img = Image.open(resource_path(f"assets/filetypes/{key[1]}.png")).convert("RGBA")
            w, h = img.size
            scale = size / max(w, h)
            ci = ctk.CTkImage(light_image=img, dark_image=img,
                              size=(int(w * scale), int(h * scale)))
        except Exception as exc:  # noqa: BLE001
            print("icon load failed", ext, exc)
            ci = None
        self._icon_cache[key] = ci
        return ci

    def _ui_icon(self, name: str, size: int = 18,
                 light: str = "#3A3436", dark: str = "#D8CFD1"):
        """Monochrome UI icon re-tinted per theme (CTkImage). Cached."""
        key = ("ui", name, size, light, dark)
        if key in self._icon_cache:
            return self._icon_cache[key]
        try:
            from PIL import Image

            base = Image.open(resource_path(f"assets/ui/{name}.png")).convert("RGBA")
            alpha = base.split()[3]

            def tint(hexcolor: str) -> "Image.Image":
                solid = Image.new("RGBA", base.size, self._hex_rgb(hexcolor) + (0,))
                solid.putalpha(alpha)
                return solid

            ci = ctk.CTkImage(light_image=tint(light), dark_image=tint(dark),
                              size=(size, size))
        except Exception as exc:  # noqa: BLE001
            print("ui icon load failed", name, exc)
            ci = None
        self._icon_cache[key] = ci
        return ci

    # ---- layout ---------------------------------------------------------- #
    def _build_sidebar(self):
        sidebar = ctk.CTkFrame(self, width=200, corner_radius=0, fg_color=SIDEBAR_FG)
        sidebar.grid(row=0, column=0, sticky="nsew")
        sidebar.grid_rowconfigure(len(NAV_ITEMS) + 1, weight=1)

        header = ctk.CTkFrame(sidebar, fg_color="transparent")
        header.grid(row=0, column=0, sticky="ew", padx=12, pady=(20, 14))
        self.sidebar_logo = self._load_logo("DashboardLogo.png", 150)
        if self.sidebar_logo is not None:
            ctk.CTkLabel(header, text="", image=self.sidebar_logo).pack(anchor="w")
        else:
            ctk.CTkLabel(
                header, text=APP_NAME, font=ctk.CTkFont(size=20, weight="bold")
            ).pack(anchor="w")
        ctk.CTkLabel(
            header, text="Convert anything", text_color=MUTED,
            font=ctk.CTkFont(size=12),
        ).pack(anchor="w", padx=(4, 0))

        self.nav_buttons: dict[str, ctk.CTkButton] = {}
        for i, name in enumerate(NAV_ITEMS, start=1):
            btn = ctk.CTkButton(
                sidebar, text="  " + name, anchor="w", height=40, corner_radius=8,
                image=self._ui_icon(NAV_ICONS[name], 18), compound="left",
                fg_color="transparent", hover_color=("#EBE0E1", "#2A2426"),
                text_color=NAV_TEXT,
                command=lambda n=name: self.show_frame(n),
            )
            btn.grid(row=i, column=0, padx=12, pady=3, sticky="ew")
            self.nav_buttons[name] = btn

        # Sun / moon appearance toggle (shows the mode you can switch to).
        self.theme_toggle = ctk.CTkButton(
            sidebar, text=SUN_GLYPH, width=44, height=36, corner_radius=18,
            font=ctk.CTkFont(family="Segoe UI Symbol", size=18),
            fg_color="transparent", border_width=1, border_color=MUTED,
            hover_color=RED_HOVER, text_color=NAV_TEXT,
            command=self._toggle_appearance,
        )
        self.theme_toggle.grid(row=len(NAV_ITEMS) + 2, column=0, padx=15, pady=20, sticky="s")

    def _build_frames(self):
        container = ctk.CTkFrame(self, corner_radius=0, fg_color="transparent")
        container.grid(row=0, column=1, sticky="nsew", padx=20, pady=20)
        container.grid_rowconfigure(0, weight=1)
        container.grid_columnconfigure(0, weight=1)

        self.frames: dict[str, ctk.CTkFrame] = {
            "Converter": self._build_converter(container),
            "Batch Convert": self._build_batch(container),
            "YouTube": self._build_youtube(container),
            "Marquee": self._build_marquee(container),
            "Home": self._build_home(container),
            "Recent": self._build_recent(container),
            "Tools": self._build_tools(container),
        }
        for frame in self.frames.values():
            frame.grid(row=0, column=0, sticky="nsew")

    def show_frame(self, name: str):
        frame = self.frames.get(name)
        if frame is None:
            return
        # grid/grid_remove (not tkraise) so a CTkScrollableFrame raises reliably
        # above the plain sibling frames sharing the same grid cell.
        for f in self.frames.values():
            if f is not frame:
                f.grid_remove()
        frame.grid()
        for n, btn in self.nav_buttons.items():
            active = n == name
            icon = (self._ui_icon(NAV_ICONS[n], 18, light="#FFFFFF", dark="#FFFFFF")
                    if active else self._ui_icon(NAV_ICONS[n], 18))
            btn.configure(
                fg_color=RED if active else "transparent",
                text_color="#FFFFFF" if active else NAV_TEXT,
                image=icon,
            )
        if name == "Recent":
            self._refresh_recent()

    def _toggle_appearance(self):
        """Flip between Dark and Light; the icon shows the mode you can switch to."""
        self.appearance_mode = "Light" if self.appearance_mode == "Dark" else "Dark"
        ctk.set_appearance_mode(self.appearance_mode)
        self.theme_toggle.configure(
            text=SUN_GLYPH if self.appearance_mode == "Dark" else MOON_GLYPH
        )

    def _section_header(self, parent, title: str, subtitle: str = ""):
        head = ctk.CTkFrame(parent, fg_color="transparent")
        head.grid(row=0, column=0, sticky="ew", padx=24, pady=(20, 8))
        ctk.CTkLabel(
            head, text=title, font=ctk.CTkFont(size=24, weight="bold"), text_color=RED
        ).pack(anchor="w")
        if subtitle:
            ctk.CTkLabel(head, text=subtitle, text_color=MUTED).pack(anchor="w", pady=(2, 0))

    def _bind_click(self, widget, command) -> None:
        """Make a whole card clickable (widget + all descendants)."""
        widget.bind("<Button-1>", lambda _e: command())
        try:
            widget.configure(cursor="hand2")
        except Exception:  # noqa: BLE001 - some widgets reject cursor
            pass
        for child in widget.winfo_children():
            self._bind_click(child, command)

    def _popular_card(self, parent, col, from_ext, to_ext, label):
        card = ctk.CTkFrame(parent, fg_color=CARD, corner_radius=12,
                            border_width=1, border_color=CARD_BORDER)
        card.grid(row=0, column=col, padx=5, sticky="ew")
        icons = ctk.CTkFrame(card, fg_color="transparent")
        icons.pack(padx=16, pady=(16, 6))
        ctk.CTkLabel(icons, text="", image=self._filetype_icon(from_ext, 30)).pack(side="left")
        ctk.CTkLabel(icons, text="", image=self._ui_icon("arrow-right", 16)).pack(side="left", padx=9)
        ctk.CTkLabel(icons, text="", image=self._filetype_icon(to_ext, 30)).pack(side="left")
        ctk.CTkLabel(
            card, text=label, text_color=TEXT, font=ctk.CTkFont(size=12, weight="bold"),
        ).pack(pady=(0, 14))
        self._bind_click(card, lambda: self.show_frame("Converter"))

    def _build_home(self, parent) -> ctk.CTkFrame:
        # NOTE: must be opaque — a transparent frame lets the stacked sibling
        # frame (same grid cell) show through when raised.
        frame = ctk.CTkScrollableFrame(parent, fg_color=("#EAE5E6", "#231F21"))
        frame.grid_columnconfigure(0, weight=1)

        # ---- hero ----
        hero = ctk.CTkFrame(frame, fg_color=SURFACE_SOFT, corner_radius=18)
        hero.grid(row=0, column=0, sticky="ew", padx=6, pady=(4, 16))
        hero.grid_columnconfigure(0, weight=1)
        text = ctk.CTkFrame(hero, fg_color="transparent")
        text.grid(row=0, column=0, sticky="w", padx=(28, 12), pady=(28, 24))
        ctk.CTkLabel(
            text, text="Convert anything,", anchor="w", text_color=TEXT,
            font=ctk.CTkFont(size=30, weight="bold"),
        ).pack(anchor="w")
        ctk.CTkLabel(
            text, text="to everything.", anchor="w", text_color=RED,
            font=ctk.CTkFont(size=30, weight="bold"),
        ).pack(anchor="w")
        ctk.CTkLabel(
            text, text="Fast, private, on-device conversion — documents, slides,\n"
                       "images, audio & video. Plus YouTube downloads.",
            anchor="w", justify="left", text_color=MUTED, font=ctk.CTkFont(size=13),
        ).pack(anchor="w", pady=(8, 16))
        actions = ctk.CTkFrame(text, fg_color="transparent")
        actions.pack(anchor="w")
        ctk.CTkButton(
            actions, text=" Convert a File", height=42, image=self._ui_icon("repeat", 18, "#FFFFFF", "#FFFFFF"),
            compound="left", font=ctk.CTkFont(size=14, weight="bold"),
            command=lambda: self.show_frame("Converter"),
        ).pack(side="left", padx=(0, 10))
        ctk.CTkButton(
            actions, text=" Batch", height=42, image=self._ui_icon("layers", 18),
            compound="left", fg_color="transparent", border_width=2, border_color=RED,
            text_color=TEXT, hover_color=("#F1DDDD", "#2E2A2C"),
            command=lambda: self.show_frame("Batch Convert"),
        ).pack(side="left", padx=(0, 10))
        ctk.CTkButton(
            actions, text=" YouTube", height=42, image=self._ui_icon("youtube", 18),
            compound="left", fg_color="transparent", border_width=2, border_color=RED,
            text_color=TEXT, hover_color=("#F1DDDD", "#2E2A2C"),
            command=lambda: self.show_frame("YouTube"),
        ).pack(side="left")

        cluster = ctk.CTkFrame(hero, fg_color="transparent")
        cluster.grid(row=0, column=1, padx=(8, 26), pady=20, sticky="e")
        for idx, ext in enumerate(["pdf", "docx", "pptx", "mp4", "jpg", "mp3"]):
            r, c = divmod(idx, 3)
            tile = ctk.CTkFrame(cluster, width=60, height=60, fg_color=CARD, corner_radius=14)
            tile.grid(row=r, column=c, padx=6, pady=6)
            tile.grid_propagate(False)
            tile.grid_rowconfigure(0, weight=1)
            tile.grid_columnconfigure(0, weight=1)
            ctk.CTkLabel(tile, text="", image=self._filetype_icon(ext, 34)).grid(row=0, column=0)

        # ---- popular conversions ----
        ctk.CTkLabel(
            frame, text="Popular conversions", anchor="w", text_color=TEXT,
            font=ctk.CTkFont(size=15, weight="bold"),
        ).grid(row=1, column=0, sticky="w", padx=10, pady=(2, 8))
        pop = ctk.CTkFrame(frame, fg_color="transparent")
        pop.grid(row=2, column=0, sticky="ew", padx=4, pady=(0, 16))
        for c in range(4):
            pop.grid_columnconfigure(c, weight=1, uniform="pop")
        for col, (fe, te, lbl) in enumerate([
            ("pdf", "docx", "PDF to Word"), ("docx", "pdf", "Word to PDF"),
            ("mp4", "mp3", "MP4 to MP3"), ("jpg", "png", "JPG to PNG"),
        ]):
            self._popular_card(pop, col, fe, te, lbl)

        # ---- supported formats ----
        ctk.CTkLabel(
            frame, text="Supported formats", anchor="w", text_color=TEXT,
            font=ctk.CTkFont(size=15, weight="bold"),
        ).grid(row=3, column=0, sticky="w", padx=10, pady=(2, 8))
        card = ctk.CTkFrame(frame, fg_color=CARD, corner_radius=12,
                            border_width=1, border_color=CARD_BORDER)
        card.grid(row=4, column=0, sticky="ew", padx=6, pady=(0, 16))
        card.grid_columnconfigure(2, weight=1)
        for r, (rep_ext, cat, fmts) in enumerate([
            ("txt", "Documents", "PDF · DOCX · TXT · MD"),
            ("pptx", "Presentations", "PPTX"),
            ("jpg", "Images", "JPG · PNG · WEBP · BMP · GIF · TIFF"),
            ("mp4", "Audio / Video", "MP4 · MP3 · WAV"),
        ]):
            ctk.CTkLabel(card, text="", image=self._filetype_icon(rep_ext, 26)).grid(
                row=r, column=0, padx=(16, 10), pady=10)
            ctk.CTkLabel(
                card, text=cat, font=ctk.CTkFont(size=12, weight="bold"),
                text_color=TEXT, width=110, anchor="w",
            ).grid(row=r, column=1, sticky="w", pady=10)
            ctk.CTkLabel(card, text=fmts, text_color=MUTED, anchor="w").grid(
                row=r, column=2, sticky="w", padx=(0, 16), pady=10)
        return frame

    def _build_tools(self, parent) -> ctk.CTkFrame:
        frame = ctk.CTkFrame(parent)
        frame.grid_columnconfigure(0, weight=1)
        self._section_header(frame, "Tools", "Status and utilities.")

        card = ctk.CTkFrame(frame, fg_color=("#DFD9DA", "#2B2629"), corner_radius=12)
        card.grid(row=1, column=0, sticky="ew", padx=24, pady=(4, 12))
        card.grid_columnconfigure(1, weight=1)

        ffmpeg_ok = shutil.which("ffmpeg") is not None
        rows = [
            ("FFmpeg (audio / video)",
             "Available" if ffmpeg_ok else "Not found on PATH",
             SUCCESS if ffmpeg_ok else "orange"),
            ("History", f"{len(self.history)} conversion(s) recorded", MUTED),
            ("Version", f"Bu D3eij {APP_VERSION}", MUTED),
        ]
        for r, (label, value, color) in enumerate(rows):
            top = 14 if r == 0 else 6
            ctk.CTkLabel(
                card, text=label, font=ctk.CTkFont(size=12, weight="bold"), anchor="w",
            ).grid(row=r, column=0, sticky="w", padx=(18, 12), pady=(top, 6))
            ctk.CTkLabel(card, text=value, text_color=color, anchor="w").grid(
                row=r, column=1, sticky="w", padx=(0, 18), pady=(top, 6)
            )

        btns = ctk.CTkFrame(frame, fg_color="transparent")
        btns.grid(row=2, column=0, sticky="w", padx=24, pady=(2, 12))
        ctk.CTkButton(
            btns, text="Open history folder", width=170,
            command=lambda: self.open_folder(HISTORY_FILE),
        ).grid(row=0, column=0, padx=(0, 10))
        ctk.CTkButton(
            btns, text="Reveal last output", width=170,
            fg_color="gray50", hover_color="gray40", command=self._open_last_output,
        ).grid(row=0, column=1)
        return frame

    def _open_last_output(self):
        for entry in self.history:
            if entry.get("ok") and entry.get("output"):
                self.open_folder(entry["output"])
                return

    # ---- recent view ----------------------------------------------------- #
    def _recent_columns(self, widget) -> None:
        """Apply the shared table column layout to a header/row frame."""
        widget.grid_columnconfigure(0, weight=1)          # File (expands)
        for col, minw in ((1, 64), (2, 64), (3, 120), (4, 110), (5, 140)):
            widget.grid_columnconfigure(col, minsize=minw)

    def _build_recent(self, parent) -> ctk.CTkFrame:
        frame = ctk.CTkFrame(parent)
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_rowconfigure(3, weight=1)
        self._section_header(frame, "Recent", "Your conversion history.")

        bar = ctk.CTkFrame(frame, fg_color="transparent")
        bar.grid(row=1, column=0, sticky="ew", padx=24)
        bar.grid_columnconfigure(0, weight=1)
        ctk.CTkButton(
            bar, text="Clear history", width=120, height=32,
            image=self._ui_icon("layers", 16), compound="left",
            fg_color=("#E7DFE0", "#2E2A2C"), hover_color=("#D9CFD1", "#3A3438"),
            text_color=TEXT, command=self.clear_history,
        ).grid(row=0, column=1, sticky="e")

        head = ctk.CTkFrame(frame, fg_color="transparent")
        head.grid(row=2, column=0, sticky="ew", padx=30, pady=(12, 0))
        self._recent_columns(head)
        for col, title in enumerate(["File", "From", "To", "Status", "Time", ""]):
            ctk.CTkLabel(
                head, text=title, text_color=MUTED, anchor="w",
                font=ctk.CTkFont(size=11, weight="bold"),
            ).grid(row=0, column=col, sticky="w", padx=(2, 6))

        self.recent_scroll = ctk.CTkScrollableFrame(frame, fg_color="transparent")
        self.recent_scroll.grid(row=3, column=0, sticky="nsew", padx=24, pady=(4, 12))
        self.recent_scroll.grid_columnconfigure(0, weight=1)
        self._refresh_recent()
        return frame

    def _refresh_recent(self):
        if not hasattr(self, "recent_scroll"):
            return
        for child in self.recent_scroll.winfo_children():
            child.destroy()
        if not self.history:
            ctk.CTkLabel(
                self.recent_scroll, text="No conversions yet — your results will appear here.",
                text_color=MUTED,
            ).grid(row=0, column=0, sticky="w", padx=10, pady=16)
            return
        for i, entry in enumerate(self.history):
            self._recent_row(i, entry)

    def _recent_row(self, i: int, entry: dict):
        src = entry.get("source", "") or ""
        out = entry.get("output", "") or ""
        ok = bool(entry.get("ok"))
        is_url = src.lower().startswith(("http://", "https://"))

        name = Path(out).name if (ok and out) else (
            self._ellipsize(src, 44) if is_url else (Path(src).name or "?"))
        from_ext = "URL" if is_url else (detect_format(src).upper() or "?")
        to_ext = (detect_format(out).upper() or "?") if out else "?"
        icon_ext = (detect_format(out) if (ok and out) else detect_format(src)) or "x"

        row = ctk.CTkFrame(self.recent_scroll, fg_color=CARD, corner_radius=8)
        row.grid(row=i, column=0, sticky="ew", pady=3, padx=2)
        self._recent_columns(row)

        fcell = ctk.CTkFrame(row, fg_color="transparent")
        fcell.grid(row=0, column=0, sticky="ew", padx=(10, 6), pady=8)
        ic = self._filetype_icon(icon_ext, 24)
        if ic is not None:
            ctk.CTkLabel(fcell, text="", image=ic).pack(side="left", padx=(0, 8))
        txt = ctk.CTkFrame(fcell, fg_color="transparent")
        txt.pack(side="left", fill="x", expand=True)
        ctk.CTkLabel(txt, text=self._ellipsize(name, 46), anchor="w", text_color=TEXT).pack(anchor="w")
        if not ok and entry.get("error"):
            ctk.CTkLabel(
                txt, text=self._ellipsize(str(entry["error"]), 60), anchor="w",
                text_color=MUTED, font=ctk.CTkFont(size=11),
            ).pack(anchor="w")

        ctk.CTkLabel(row, text=from_ext, text_color=MUTED, anchor="w").grid(row=0, column=1, sticky="w")
        ctk.CTkLabel(row, text=to_ext, text_color=TEXT, anchor="w").grid(row=0, column=2, sticky="w")
        ctk.CTkLabel(
            row, text=("✓ Completed" if ok else "✕ Failed"),
            text_color=(SUCCESS if ok else ERROR), anchor="w",
        ).grid(row=0, column=3, sticky="w")
        ctk.CTkLabel(row, text=entry.get("time", ""), text_color=MUTED, anchor="w").grid(
            row=0, column=4, sticky="w")

        if ok and out:
            # Only build the actions frame when it has buttons — an empty
            # CTkFrame keeps its default 200px height and stretches the row.
            act = ctk.CTkFrame(row, fg_color="transparent")
            act.grid(row=0, column=5, sticky="e", padx=(0, 8))
            ctk.CTkButton(
                act, text="Open", width=50, height=26,
                command=lambda p=out: self.open_path(p),
            ).pack(side="left", padx=(0, 6))
            ctk.CTkButton(
                act, text="Folder", width=60, height=26,
                fg_color="transparent", border_width=1, border_color=CARD_BORDER,
                text_color=TEXT, hover_color=("#EBE0E1", "#332D30"),
                command=lambda p=out: self.open_folder(p),
            ).pack(side="left")

    def add_history(self, src, out, ok: bool, error="") -> None:
        entry = {
            "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "source": str(src) if src else "",
            "output": str(out) if out else "",
            "ok": bool(ok),
            "error": str(error),
        }
        self.history.insert(0, entry)
        del self.history[MAX_HISTORY:]
        save_history(self.history)
        self.after(0, self._refresh_recent)

    def clear_history(self):
        self.history = []
        save_history(self.history)
        self._refresh_recent()

    def open_path(self, path):
        try:
            os.startfile(str(path))  # type: ignore[attr-defined]
        except Exception as exc:  # noqa: BLE001
            print("Open failed:", exc)

    def open_folder(self, path):
        p = Path(path)
        try:
            if p.exists():
                subprocess.Popen(["explorer", "/select,", str(p)])
            else:
                os.startfile(str(p.parent))  # type: ignore[attr-defined]
        except Exception as exc:  # noqa: BLE001
            print("Open folder failed:", exc)

    # ---- converter view -------------------------------------------------- #
    def _build_converter(self, parent) -> ctk.CTkFrame:
        frame = ctk.CTkFrame(parent)
        frame.grid_columnconfigure(0, weight=1)

        self._section_header(frame, "Converter", "Drop a file, choose a format, convert.")

        # ---- drop zone ----
        self.drop_zone = ctk.CTkFrame(
            frame, height=200, corner_radius=14, border_width=2, border_color=DROP_BORDER
        )
        self.drop_zone.grid(row=1, column=0, sticky="ew", padx=24, pady=(4, 14))
        self.drop_zone.grid_propagate(False)
        self.drop_zone.grid_columnconfigure(0, weight=1)
        self.drop_zone.grid_rowconfigure(0, weight=1)
        inner = ctk.CTkFrame(self.drop_zone, fg_color="transparent")
        inner.grid(row=0, column=0)
        # On-theme folder icon from bundled assets (renders in the frozen exe;
        # AppLogo.png was never bundled). Swapped for the file's type icon on drop.
        self.drop_icon = self._ui_icon("folder-open", 46, light=RED, dark=RED_BRIGHT)
        self.drop_icon_label = ctk.CTkLabel(inner, text="", image=self.drop_icon)
        self.drop_icon_label.pack(pady=(0, 8))
        self.drop_primary = ctk.CTkLabel(
            inner, text="Drag & drop a file here",
            font=ctk.CTkFont(size=16, weight="bold"),
        )
        self.drop_primary.pack()
        self.drop_secondary = ctk.CTkLabel(inner, text="or click to browse", text_color=MUTED)
        self.drop_secondary.pack(pady=(2, 0))
        self._register_drop(
            self.drop_zone,
            (self.drop_zone, inner, self.drop_icon_label, self.drop_primary, self.drop_secondary),
            self.on_drop, self.browse_file,
        )
        # Wrap the filename / hint text to the zone width so long names stay inside.
        self._make_labels_wrap(self.drop_zone, (self.drop_primary, self.drop_secondary))

        # ---- controls card ----
        card = ctk.CTkFrame(frame, fg_color=("#DFD9DA", "#2B2629"), corner_radius=12)
        card.grid(row=2, column=0, sticky="ew", padx=24, pady=(0, 12))
        card.grid_columnconfigure(0, weight=1)

        fmt = ctk.CTkFrame(card, fg_color="transparent")
        fmt.grid(row=0, column=0, sticky="w", padx=18, pady=(16, 8))
        ctk.CTkLabel(fmt, text="CONVERT FROM", font=ctk.CTkFont(size=11, weight="bold"),
                     text_color=MUTED).grid(row=0, column=0, sticky="w", padx=2)
        ctk.CTkLabel(fmt, text="CONVERT TO", font=ctk.CTkFont(size=11, weight="bold"),
                     text_color=MUTED).grid(row=0, column=2, sticky="w", padx=(16, 0))
        self.from_menu = ctk.CTkOptionMenu(fmt, values=["-"], state="disabled", width=150)
        self.from_menu.grid(row=1, column=0, pady=(2, 0))
        ctk.CTkLabel(fmt, text="→", font=ctk.CTkFont(size=22, weight="bold"),
                     text_color=RED).grid(row=1, column=1, padx=14)
        self.to_menu = ctk.CTkOptionMenu(fmt, values=["-"], state="disabled", width=150)
        self.to_menu.grid(row=1, column=2, pady=(2, 0))

        exp = ctk.CTkFrame(card, fg_color="transparent")
        exp.grid(row=1, column=0, sticky="ew", padx=18, pady=(6, 4))
        exp.grid_columnconfigure(0, weight=1)
        self.export_label = ctk.CTkLabel(
            exp, text="Output: next to the source file", text_color=MUTED, anchor="w",
        )
        self.export_label.grid(row=0, column=0, sticky="ew", padx=(2, 8))
        self.export_btn = ctk.CTkButton(
            exp, text="Choose Folder", width=130, command=self.choose_export_path
        )
        self.export_btn.grid(row=0, column=1, padx=(0, 8))
        self.clear_btn = ctk.CTkButton(
            exp, text="Clear", width=80, fg_color="gray50",
            hover_color="gray40", command=self.clear_converter,
        )
        self.clear_btn.grid(row=0, column=2)

        self.convert_btn = GradientButton(
            card, text="Convert Now", height=46,
            command=self.on_convert_click, state="disabled",
        )
        self.convert_btn.grid(row=2, column=0, sticky="ew", padx=18, pady=(8, 16))

        self.progress = ctk.CTkProgressBar(frame, height=8)
        self.progress.set(0)
        self.progress.grid(row=3, column=0, sticky="ew", padx=24, pady=(0, 4))
        self.progress.grid_remove()  # only shown while converting

        self.status = ctk.CTkLabel(
            frame, text="Drop a file to begin.", text_color=MUTED,
            wraplength=620, justify="left", anchor="w",
        )
        self.status.grid(row=4, column=0, sticky="w", padx=24, pady=(4, 16))
        return frame

    # ---- batch view ------------------------------------------------------ #
    def _build_batch(self, parent) -> ctk.CTkFrame:
        frame = ctk.CTkFrame(parent)
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_rowconfigure(3, weight=1)

        self._section_header(frame, "Batch Convert", "Convert many files to one format at once.")

        self.batch_drop = ctk.CTkFrame(
            frame, height=120, corner_radius=14, border_width=2, border_color=DROP_BORDER
        )
        self.batch_drop.grid(row=1, column=0, sticky="ew", padx=24, pady=(4, 12))
        self.batch_drop.grid_propagate(False)
        self.batch_drop.grid_columnconfigure(0, weight=1)
        self.batch_drop.grid_rowconfigure(0, weight=1)
        binner = ctk.CTkFrame(self.batch_drop, fg_color="transparent")
        binner.grid(row=0, column=0)
        self.batch_primary = ctk.CTkLabel(
            binner, text="Drop multiple files here",
            font=ctk.CTkFont(size=15, weight="bold"),
        )
        self.batch_primary.pack()
        self.batch_label = ctk.CTkLabel(binner, text="or click to browse", text_color=MUTED)
        self.batch_label.pack(pady=(2, 0))
        self._register_drop(
            self.batch_drop, (self.batch_drop, binner, self.batch_primary, self.batch_label),
            self.on_batch_drop, self.browse_batch,
        )

        ctrl = ctk.CTkFrame(frame, fg_color="transparent")
        ctrl.grid(row=2, column=0, sticky="w", padx=24, pady=8)
        ctk.CTkLabel(ctrl, text="Convert all to:").grid(row=0, column=0, padx=(0, 8))
        self.batch_to = ctk.CTkOptionMenu(ctrl, values=["-"], width=140, state="disabled")
        self.batch_to.grid(row=0, column=1, padx=(0, 16))
        self.batch_btn = ctk.CTkButton(
            ctrl, text=" Convert All", width=140, font=ctk.CTkFont(weight="bold"),
            image=self._ui_icon("repeat", 16, "#FFFFFF", "#FFFFFF"), compound="left",
            command=self.on_batch_convert, state="disabled",
        )
        self.batch_btn.grid(row=0, column=2)

        self.batch_list = ctk.CTkTextbox(frame, corner_radius=10)
        self.batch_list.grid(row=3, column=0, sticky="nsew", padx=24, pady=8)

        self.batch_progress = ctk.CTkProgressBar(frame, height=8)
        self.batch_progress.set(0)
        self.batch_progress.grid(row=4, column=0, sticky="ew", padx=24, pady=(0, 18))
        return frame

    # ---- drag & drop helpers -------------------------------------------- #
    def _register_drop(self, zone, widgets, on_drop, on_click):
        def on_enter(_e):
            zone.configure(border_color=RED_BRIGHT)

        def on_leave(_e):
            zone.configure(border_color=DROP_BORDER)

        def on_drop_wrapped(event):
            on_leave(event)
            on_drop(event)

        for widget in widgets:
            try:
                widget.drop_target_register(DND_FILES)
                widget.dnd_bind("<<Drop>>", on_drop_wrapped)
                widget.dnd_bind("<<DropEnter>>", on_enter)
                widget.dnd_bind("<<DropLeave>>", on_leave)
            except Exception as exc:  # noqa: BLE001
                print("Drag & drop registration failed:", exc)
            widget.bind("<Button-1>", lambda _e: on_click())
            try:
                widget.configure(cursor="hand2")
            except Exception:  # noqa: BLE001
                pass

    def _parse_drop(self, event) -> list[Path]:
        try:
            paths = self.tk.splitlist(event.data)
        except Exception:  # noqa: BLE001
            paths = [event.data]
        return [Path(p) for p in paths]

    def _make_labels_wrap(self, container, labels, margin: int = 48) -> None:
        """Keep `labels` wrapped within `container`'s width so long text (e.g. a
        long filename) can never spill past the container's border."""
        def _update(event):
            wrap = max(event.width - margin, 120)
            for lbl in labels:
                lbl.configure(wraplength=wrap)

        container.bind("<Configure>", _update)

    @staticmethod
    def _ellipsize(text: str, limit: int = 52) -> str:
        """Shorten `text` with a middle ellipsis (keeps start and end visible)."""
        if len(text) <= limit:
            return text
        head = (limit - 1) // 2
        tail = limit - 1 - head
        return f"{text[:head]}…{text[-tail:]}"

    # ---- converter actions ---------------------------------------------- #
    def on_drop(self, event):
        paths = self._parse_drop(event)
        if paths:
            self.set_file(paths[0])

    def browse_file(self):
        path = filedialog.askopenfilename(title="Select a file to convert")
        if path:
            self.set_file(Path(path))

    def set_file(self, path: Path):
        if not path.is_file():
            self.status.configure(text="Please drop a single file.", text_color="orange")
            return
        self.selected_file = path
        ext = detect_format(path)
        targets = compatible_targets(ext)
        try:
            size = human_size(path.stat().st_size)
        except OSError:
            size = "?"
        self.drop_primary.configure(text=path.name)
        ft_icon = self._filetype_icon(ext, 52)
        if ft_icon is not None:
            self.drop_icon_label.configure(image=ft_icon)
        self.from_menu.configure(values=[ext or "?"])
        self.from_menu.set(ext or "?")
        if targets:
            self.drop_secondary.configure(
                text=f"{category_of(ext)}  ·  {size}  ·  click to choose another")
            self.to_menu.configure(values=targets, state="normal")
            self.to_menu.set(targets[0])
            self.convert_btn.configure(state="normal")
            self.status.configure(text=f"Ready to convert {path.name}", text_color=MUTED)
        else:
            self.drop_secondary.configure(text="Unsupported format  ·  click to choose another")
            self.to_menu.configure(values=["-"], state="disabled")
            self.to_menu.set("-")
            self.convert_btn.configure(state="disabled")
            self.status.configure(text=f"Unsupported format: .{ext}", text_color="orange")

    def choose_export_path(self):
        folder = filedialog.askdirectory(title="Choose where to save converted files")
        if folder:
            self.export_dir = Path(folder)
            self.export_label.configure(text=f"Output: {self._ellipsize(str(self.export_dir))}")

    def clear_converter(self):
        """Reset the converter to a clean slate for the next file."""
        self.selected_file = None
        self.export_dir = None
        if self.drop_icon is not None:
            self.drop_icon_label.configure(image=self.drop_icon)
        self.drop_primary.configure(text="Drag & drop a file here")
        self.drop_secondary.configure(text="or click to browse")
        self.from_menu.configure(values=["-"], state="disabled")
        self.from_menu.set("-")
        self.to_menu.configure(values=["-"], state="disabled")
        self.to_menu.set("-")
        self.convert_btn.configure(state="disabled")
        self.progress.stop()
        self.progress.grid_remove()
        self.progress.configure(mode="determinate")
        self.progress.set(0)
        self.export_label.configure(text="Output: next to the source file")
        self.status.configure(text="Drop a file to begin.", text_color=MUTED)

    def on_convert_click(self):
        if not self.selected_file:
            return
        target = self.to_menu.get()
        if target in ("-", ""):
            return
        self.convert_btn.configure(state="disabled")
        self.convert_btn.start_busy()
        self.progress.grid()
        self.progress.configure(mode="indeterminate")
        self.progress.start()
        self.status.configure(
            text=f"Converting {self.selected_file.name} to .{target} …",
            text_color=MUTED,
        )
        threading.Thread(
            target=self._convert_worker, args=(self.selected_file, target), daemon=True
        ).start()

    def _convert_worker(self, src: Path, target: str):
        try:
            out = convert_file(src, target, self.export_dir)
            self.after(0, self._convert_done, src, out, None)
        except Exception as exc:  # noqa: BLE001
            traceback.print_exc()
            self.after(0, self._convert_done, src, None, exc)

    def _convert_done(self, src: Path, out: Path | None, error: Exception | None):
        self.progress.stop()
        self.progress.grid_remove()
        self.progress.configure(mode="determinate")
        self.progress.set(0)
        self.convert_btn.configure(state="normal")
        self.convert_btn.stop_busy()
        if error:
            self.status.configure(text=f"✕  {error}", text_color=ERROR)
            self.add_history(src, None, False, error)
        else:
            self.status.configure(text=f"✓  Saved to {out}", text_color=SUCCESS)
            self.add_history(src, out, True)

    # ---- batch actions --------------------------------------------------- #
    def on_batch_drop(self, event):
        self.set_batch(self._parse_drop(event))

    def browse_batch(self):
        paths = filedialog.askopenfilenames(title="Select files to convert")
        if paths:
            self.set_batch([Path(p) for p in paths])

    def set_batch(self, paths: list[Path]):
        self.batch_files = [p for p in paths if p.is_file()]
        if not self.batch_files:
            return
        common: set[str] | None = None
        for p in self.batch_files:
            targets = set(compatible_targets(detect_format(p)))
            common = targets if common is None else (common & targets)
        common_sorted = sorted(common) if common else []

        self.batch_primary.configure(text=f"{len(self.batch_files)} file(s) selected")
        self.batch_label.configure(text="click to choose different files")
        self.batch_list.delete("1.0", "end")
        for p in self.batch_files:
            self.batch_list.insert("end", f"- {p.name}\n")

        if common_sorted:
            self.batch_to.configure(values=common_sorted, state="normal")
            self.batch_to.set(common_sorted[0])
            self.batch_btn.configure(state="normal")
        else:
            self.batch_to.configure(values=["-"], state="disabled")
            self.batch_to.set("-")
            self.batch_btn.configure(state="disabled")
            self.batch_list.insert(
                "end", "\nNo common target format. Try files of the same type.\n"
            )

    def on_batch_convert(self):
        target = self.batch_to.get()
        if target in ("-", "") or not self.batch_files:
            return
        self.batch_btn.configure(state="disabled")
        self.batch_progress.set(0)
        self.batch_list.delete("1.0", "end")
        files = list(self.batch_files)
        threading.Thread(
            target=self._batch_worker, args=(files, target), daemon=True
        ).start()

    def _batch_worker(self, files: list[Path], target: str):
        total = len(files)
        for i, p in enumerate(files, start=1):
            try:
                out = convert_file(p, target)
                self.after(0, self._batch_log, f"OK   {p.name} -> {out.name}")
                # Marshal history updates to the main thread (self.history is shared
                # with the UI; mutating it off-thread can race with _refresh_recent).
                self.after(0, self.add_history, p, out, True)
            except Exception as exc:  # noqa: BLE001
                traceback.print_exc()
                self.after(0, self._batch_log, f"FAIL {p.name}: {exc}")
                self.after(0, self.add_history, p, None, False, exc)
            self.after(0, self.batch_progress.set, i / total)
        self.after(0, lambda: self.batch_btn.configure(state="normal"))
        self.after(0, self._batch_log, f"\nDone: {total} file(s) processed.")

    def _batch_log(self, msg: str):
        self.batch_list.insert("end", msg + "\n")
        self.batch_list.see("end")

    # ---- youtube view ---------------------------------------------------- #
    def _build_youtube(self, parent) -> ctk.CTkFrame:
        frame = ctk.CTkFrame(parent)
        frame.grid_columnconfigure(0, weight=1)
        self._section_header(frame, "YouTube", "Paste a link, pick a format, download.")

        card = ctk.CTkFrame(frame, fg_color=("#DFD9DA", "#2B2629"), corner_radius=12)
        card.grid(row=1, column=0, sticky="ew", padx=24, pady=(4, 12))
        card.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(
            card, text="VIDEO URL", font=ctk.CTkFont(size=11, weight="bold"), text_color=MUTED,
        ).grid(row=0, column=0, sticky="w", padx=18, pady=(16, 2))
        self.yt_url = ctk.CTkEntry(
            card, placeholder_text="https://www.youtube.com/watch?v=…", height=40,
        )
        self.yt_url.grid(row=1, column=0, sticky="ew", padx=18, pady=(0, 10))
        self.yt_url.bind("<Return>", lambda _e: self.on_youtube_download())

        opts = ctk.CTkFrame(card, fg_color="transparent")
        opts.grid(row=2, column=0, sticky="w", padx=18, pady=(0, 6))
        ctk.CTkLabel(opts, text="Format:").grid(row=0, column=0, padx=(0, 10))
        self.yt_format = ctk.CTkSegmentedButton(
            opts, values=["MP4", "MP3"], selected_color=RED, selected_hover_color=RED_HOVER,
        )
        self.yt_format.set("MP4")
        self.yt_format.grid(row=0, column=1)
        ctk.CTkLabel(
            opts, text="MP4 = video · MP3 = audio (192 kbps)", text_color=MUTED,
        ).grid(row=0, column=2, padx=(14, 0))

        self.yt_btn = ctk.CTkButton(
            card, text=" Download", height=46, font=ctk.CTkFont(size=15, weight="bold"),
            image=self._ui_icon("download", 18, "#FFFFFF", "#FFFFFF"), compound="left",
            command=self.on_youtube_download,
        )
        self.yt_btn.grid(row=3, column=0, sticky="ew", padx=18, pady=(8, 16))

        self.yt_progress = ctk.CTkProgressBar(frame, height=8)
        self.yt_progress.set(0)
        self.yt_progress.grid(row=2, column=0, sticky="ew", padx=24, pady=(0, 4))
        self.yt_progress.grid_remove()

        self.yt_status = ctk.CTkLabel(
            frame, text="Paste a video link to begin.", text_color=MUTED,
            wraplength=620, justify="left", anchor="w",
        )
        self.yt_status.grid(row=3, column=0, sticky="w", padx=24, pady=(4, 16))
        return frame

    def on_youtube_download(self):
        url = self.yt_url.get().strip()
        if not url:
            self.yt_status.configure(text="Please paste a video URL.", text_color="orange")
            return
        fmt = self.yt_format.get().lower()
        folder = filedialog.askdirectory(title="Choose where to save the download")
        if not folder:
            self.yt_status.configure(text="Download cancelled (no folder chosen).", text_color=MUTED)
            return
        self.yt_btn.configure(state="disabled")
        self.yt_progress.grid()
        self.yt_progress.configure(mode="determinate")
        self.yt_progress.set(0)
        self.yt_status.configure(text=f"Starting {fmt.upper()} download…", text_color=MUTED)
        threading.Thread(
            target=self._youtube_worker, args=(url, fmt, folder), daemon=True
        ).start()

    def _youtube_worker(self, url: str, fmt: str, folder: str):
        try:
            out = download_youtube(url, fmt, folder, progress_hook=self._yt_hook)
            self.after(0, self._youtube_done, url, out, None)
        except Exception as exc:  # noqa: BLE001
            traceback.print_exc()
            self.after(0, self._youtube_done, url, None, exc)

    def _yt_hook(self, d: dict):
        # Runs on the worker thread; marshal everything to the UI thread.
        status = d.get("status")
        if status == "downloading":
            total = d.get("total_bytes") or d.get("total_bytes_estimate")
            done = d.get("downloaded_bytes", 0)
            if total:
                self.after(0, self._yt_progress_set, done / total,
                           f"Downloading… {int(done / total * 100)}%")
            else:
                self.after(0, self._yt_progress_set, None, "Downloading…")
        elif status == "finished":
            self.after(0, self._yt_progress_set, 1.0, "Processing with ffmpeg…")

    def _yt_progress_set(self, frac, text):
        if frac is None:
            if self.yt_progress.cget("mode") != "indeterminate":
                self.yt_progress.configure(mode="indeterminate")
                self.yt_progress.start()
        else:
            if self.yt_progress.cget("mode") != "determinate":
                self.yt_progress.stop()
                self.yt_progress.configure(mode="determinate")
            self.yt_progress.set(frac)
        if text:
            self.yt_status.configure(text=text, text_color=MUTED)

    def _youtube_done(self, url: str, out: Path | None, error: Exception | None):
        self.yt_progress.stop()
        self.yt_progress.grid_remove()
        self.yt_progress.configure(mode="determinate")
        self.yt_progress.set(0)
        self.yt_btn.configure(state="normal")
        if error:
            self.yt_status.configure(text=f"✕  {error}", text_color=ERROR)
            self.add_history(url, None, False, error)
        else:
            self.yt_status.configure(text=f"✓  Saved to {out}", text_color=SUCCESS)
            self.add_history(url, out, True)

    # ---- marquee (image editing) view ------------------------------------ #
    def _build_marquee(self, parent) -> ctk.CTkFrame:
        frame = ctk.CTkFrame(parent)
        frame.grid_columnconfigure(0, weight=1)
        self._section_header(
            frame, "Marquee",
            "Remove an image background — export a clean transparent PNG.",
        )

        # ---- drop zone ----
        self.mq_drop = ctk.CTkFrame(
            frame, height=200, corner_radius=14, border_width=2, border_color=DROP_BORDER
        )
        self.mq_drop.grid(row=1, column=0, sticky="ew", padx=24, pady=(4, 14))
        self.mq_drop.grid_propagate(False)
        self.mq_drop.grid_columnconfigure(0, weight=1)
        self.mq_drop.grid_rowconfigure(0, weight=1)
        inner = ctk.CTkFrame(self.mq_drop, fg_color="transparent")
        inner.grid(row=0, column=0)
        self.mq_drop_icon = self._ui_icon("sparkles", 46, light=RED, dark=RED_BRIGHT)
        self.mq_icon_label = ctk.CTkLabel(inner, text="", image=self.mq_drop_icon)
        self.mq_icon_label.pack(pady=(0, 8))
        self.mq_primary = ctk.CTkLabel(
            inner, text="Drag & drop an image here",
            font=ctk.CTkFont(size=16, weight="bold"),
        )
        self.mq_primary.pack()
        self.mq_secondary = ctk.CTkLabel(inner, text="or click to browse", text_color=MUTED)
        self.mq_secondary.pack(pady=(2, 0))
        self._register_drop(
            self.mq_drop,
            (self.mq_drop, inner, self.mq_icon_label, self.mq_primary, self.mq_secondary),
            self.on_marquee_drop, self.browse_marquee,
        )
        self._make_labels_wrap(self.mq_drop, (self.mq_primary, self.mq_secondary))

        # ---- controls card ----
        card = ctk.CTkFrame(frame, fg_color=("#DFD9DA", "#2B2629"), corner_radius=12)
        card.grid(row=2, column=0, sticky="ew", padx=24, pady=(0, 12))
        card.grid_columnconfigure(0, weight=1)
        ctk.CTkLabel(
            card, text="Drops the background and keeps your subject on a fully "
                       "transparent canvas. Output is always a PNG.",
            text_color=MUTED, anchor="w", justify="left", wraplength=620,
        ).grid(row=0, column=0, sticky="ew", padx=18, pady=(16, 8))

        # ---- quality (model tier) selector ----
        qual = ctk.CTkFrame(card, fg_color="transparent")
        qual.grid(row=1, column=0, sticky="w", padx=18, pady=(0, 2))
        ctk.CTkLabel(qual, text="QUALITY", font=ctk.CTkFont(size=11, weight="bold"),
                     text_color=MUTED).grid(row=0, column=0, padx=(0, 12))
        self.mq_model = ctk.CTkSegmentedButton(
            qual, values=list(BG_MODELS), command=self._on_mq_model_change,
            selected_color=RED, selected_hover_color=RED_HOVER,
        )
        self.mq_model.set(DEFAULT_BG_TIER)
        self.mq_model.grid(row=0, column=1)
        self.mq_model_caption = ctk.CTkLabel(
            card, text=BG_MODELS[DEFAULT_BG_TIER][1], text_color=MUTED,
            anchor="w", justify="left", wraplength=620,
        )
        self.mq_model_caption.grid(row=2, column=0, sticky="ew", padx=18, pady=(0, 10))

        self.mq_btn = ctk.CTkButton(
            card, text=" Remove Background", height=46,
            font=ctk.CTkFont(size=15, weight="bold"),
            image=self._ui_icon("sparkles", 18, "#FFFFFF", "#FFFFFF"), compound="left",
            command=self.on_marquee_remove, state="disabled",
        )
        self.mq_btn.grid(row=3, column=0, sticky="ew", padx=18, pady=(0, 16))

        self.mq_progress = ctk.CTkProgressBar(frame, height=8)
        self.mq_progress.set(0)
        self.mq_progress.grid(row=3, column=0, sticky="ew", padx=24, pady=(0, 4))
        self.mq_progress.grid_remove()  # only shown while processing

        self.mq_status = ctk.CTkLabel(
            frame, text="Drop an image to begin.", text_color=MUTED,
            wraplength=620, justify="left", anchor="w",
        )
        self.mq_status.grid(row=4, column=0, sticky="w", padx=24, pady=(4, 16))
        return frame

    # ---- marquee actions ------------------------------------------------- #
    def _on_mq_model_change(self, tier: str):
        blurb = BG_MODELS.get(tier, ("", ""))[1]
        self.mq_model_caption.configure(text=blurb)

    def on_marquee_drop(self, event):
        paths = self._parse_drop(event)
        if paths:
            self.set_marquee_file(paths[0])

    def browse_marquee(self):
        path = filedialog.askopenfilename(
            title="Select an image",
            filetypes=[("Images", "*.png *.jpg *.jpeg *.webp *.bmp *.gif *.tiff"),
                       ("All files", "*.*")],
        )
        if path:
            self.set_marquee_file(Path(path))

    def set_marquee_file(self, path: Path):
        if not path.is_file():
            self.mq_status.configure(text="Please drop a single image.", text_color="orange")
            return
        ext = detect_format(path)
        self.mq_primary.configure(text=path.name)
        if ext not in IMAGE_EXTS:
            self.marquee_file = None
            if self.mq_drop_icon is not None:
                self.mq_icon_label.configure(image=self.mq_drop_icon)
            self.mq_secondary.configure(text="Not an image  ·  click to choose another")
            self.mq_btn.configure(state="disabled")
            self.mq_status.configure(
                text=f"Unsupported file: .{ext or '?'} — pick an image.", text_color="orange")
            return
        self.marquee_file = path
        try:
            size = human_size(path.stat().st_size)
        except OSError:
            size = "?"
        ft_icon = self._filetype_icon(ext, 52)
        if ft_icon is not None:
            self.mq_icon_label.configure(image=ft_icon)
        self.mq_secondary.configure(text=f"Image  ·  {size}  ·  click to choose another")
        self.mq_btn.configure(state="normal")
        self.mq_status.configure(
            text=f"Ready to remove the background from {path.name}", text_color=MUTED)

    def on_marquee_remove(self):
        src = self.marquee_file
        if not src:
            return
        out = filedialog.asksaveasfilename(
            title="Save transparent PNG as",
            defaultextension=".png",
            initialfile=f"{src.stem}_no-bg.png",
            filetypes=[("PNG image", "*.png")],
        )
        if not out:
            self.mq_status.configure(
                text="Cancelled (no save location chosen).", text_color=MUTED)
            return
        tier = self.mq_model.get() or DEFAULT_BG_TIER
        model = BG_MODELS.get(tier, BG_MODELS[DEFAULT_BG_TIER])[0]
        self.mq_btn.configure(state="disabled")
        self.mq_progress.grid()
        self.mq_progress.configure(mode="indeterminate")
        self.mq_progress.start()
        self.mq_status.configure(
            text=f"Removing background with {tier}… "
                 "(the first use of a model downloads it once)",
            text_color=MUTED,
        )
        threading.Thread(
            target=self._marquee_worker, args=(src, Path(out), model), daemon=True
        ).start()

    def _marquee_worker(self, src: Path, out: Path, model: str):
        try:
            result = remove_background(src, out, model)
            self.after(0, self._marquee_done, src, result, None)
        except Exception as exc:  # noqa: BLE001
            traceback.print_exc()
            self.after(0, self._marquee_done, src, None, exc)

    def _marquee_done(self, src: Path, out: Path | None, error: Exception | None):
        self.mq_progress.stop()
        self.mq_progress.grid_remove()
        self.mq_progress.configure(mode="determinate")
        self.mq_progress.set(0)
        self.mq_btn.configure(state="normal")
        if error:
            self.mq_status.configure(text=f"✕  {error}", text_color=ERROR)
            self.add_history(src, None, False, error)
        else:
            self.mq_status.configure(text=f"✓  Saved to {out}", text_color=SUCCESS)
            self.add_history(src, out, True)


def _run_cli(args) -> int:
    """Headless mode: --convert, --download, or --remove-bg. Returns an exit code."""
    import argparse

    parser = argparse.ArgumentParser(prog=APP_NAME, description="Bu D3eij file converter")
    parser.add_argument(
        "--convert", nargs=2, metavar=("FILE", "FORMAT"),
        help="Convert FILE to FORMAT (e.g. --convert photo.png jpg) and exit.",
    )
    parser.add_argument(
        "--download", nargs=2, metavar=("URL", "FORMAT"),
        help="Download URL as FORMAT (mp3/mp4) into the current folder and exit.",
    )
    parser.add_argument(
        "--remove-bg", metavar="FILE",
        help="Remove FILE's background, save a transparent PNG next to it, and exit.",
    )
    ns = parser.parse_args(args)
    if ns.convert:
        src, target = ns.convert
        try:
            out = convert_file(src, target)
            print(f"Saved: {out}")
            return 0
        except Exception as exc:  # noqa: BLE001
            print(f"Error: {exc}", file=sys.stderr)
            return 1
    if ns.download:
        url, fmt = ns.download
        try:
            out = download_youtube(url, fmt, os.getcwd())
            print(f"Saved: {out}")
            return 0
        except Exception as exc:  # noqa: BLE001
            print(f"Error: {exc}", file=sys.stderr)
            return 1
    if ns.remove_bg:
        try:
            out = remove_background(ns.remove_bg)
            print(f"Saved: {out}")
            return 0
        except Exception as exc:  # noqa: BLE001
            print(f"Error: {exc}", file=sys.stderr)
            return 1
    return 0


def main():
    # If launched with flags (e.g. --convert), run headless; otherwise open the GUI.
    if len(sys.argv) > 1 and sys.argv[1].startswith("-"):
        sys.exit(_run_cli(sys.argv[1:]))
    App().mainloop()


if __name__ == "__main__":
    main()
