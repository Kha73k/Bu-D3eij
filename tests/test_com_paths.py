"""Verify the COM paths: direct-Word DOCX->PDF and owned-PowerPoint PPTX->PDF,
plus the Word-leak fix (a failing conversion must not leave WINWORD.EXE).

Run with the venv python: .\\.venv\\Scripts\\python tests\\test_com_paths.py
Needs MS Word + PowerPoint installed.
"""
import subprocess
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import app  # noqa: E402
from bud3eij.converters import _docx_to_pdf_word  # noqa: E402

PASS, FAIL = 0, 0


def check(name, cond, detail=""):
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  OK   {name}")
    else:
        FAIL += 1
        print(f"  FAIL {name}  {detail}")


def winword_count() -> int:
    out = subprocess.run(
        ["tasklist", "/FI", "IMAGENAME eq WINWORD.EXE", "/FO", "CSV", "/NH"],
        capture_output=True, text=True).stdout
    return out.count("WINWORD.EXE")


tmp = Path(tempfile.mkdtemp(prefix="comcheck_"))

# build a real docx
import docx  # noqa: E402

d = docx.Document()
d.add_heading("COM path test", level=1)
d.add_paragraph("This paragraph proves the direct-COM Word conversion works.")
src_docx = tmp / "test.docx"
d.save(str(src_docx))

before = winword_count()
out = app.convert_file(src_docx, "pdf")
check("docx->pdf via direct Word COM", out.exists() and out.read_bytes()[:4] == b"%PDF")
check("no Word process leaked on success", winword_count() == before,
      f"before={before} after={winword_count()}")

# failure path must not leak a hidden WINWORD.EXE either
bad = tmp / "broken.docx"
bad.write_text("not a real docx", encoding="utf-8")
before = winword_count()
try:
    _docx_to_pdf_word(bad, tmp / "broken.pdf")
    check("broken docx raises", False)
except Exception:  # noqa: BLE001 - any COM error is fine, we test the cleanup
    check("broken docx raises", True)
check("no Word process leaked on failure", winword_count() == before,
      f"before={before} after={winword_count()}")

# pptx -> pdf (PowerPoint not running -> we own the instance -> Quit allowed)
from pptx import Presentation  # noqa: E402

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "COM ownership test"
src_pptx = tmp / "deck.pptx"
prs.save(str(src_pptx))
out = app.convert_file(src_pptx, "pdf")
check("pptx->pdf via owned PowerPoint", out.exists() and out.read_bytes()[:4] == b"%PDF")

print(f"\n==== {PASS} passed, {FAIL} failed ====")
raise SystemExit(1 if FAIL else 0)
